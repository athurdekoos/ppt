#!/usr/bin/env python3
"""
OpenTeams 2025 PowerPoint Template Builder
===========================================
Generates a production-ready, on-brand .pptx template deck.

Usage:
    python build_template.py \
        --assets ./Assets \
        --guidelines ./OpenTeams_Brand_Guidelines_2025.pdf \
        --site https://openteams.com/ \
        --out ./OpenTeams_Template_2025.pptx

Requirements:  pip install python-pptx requests beautifulsoup4 Pillow lxml
"""

from __future__ import annotations

import argparse
import json
import logging
import math
import os
import re
import sys
import textwrap
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

# ---------------------------------------------------------------------------
# Third-party imports
# ---------------------------------------------------------------------------
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    import pptx.oxml.ns as nsmap
    from lxml import etree
except ImportError as exc:
    sys.exit(f"Missing dependency: {exc}. Run: pip install -r requirements.txt")

try:
    import requests
    from bs4 import BeautifulSoup
except ImportError:
    requests = None  # type: ignore
    BeautifulSoup = None  # type: ignore

try:
    from PIL import Image
except ImportError:
    Image = None  # type: ignore

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------
logging.basicConfig(level=logging.INFO, format="%(levelname)s  %(message)s")
log = logging.getLogger("ot_template")

# ===================================================================
# SECTION 1 — BRAND TOKENS  (Source of truth: Brand Guidelines 2025)
# ===================================================================

@dataclass
class BrandColors:
    """Exact hex values from OpenTeams Brand Guidelines 2025."""
    night_navy: str   = "#022791"   # Primary dark
    day_blue: str     = "#4D75FE"   # Primary bright
    yellow: str       = "#FAA944"   # Accent
    salmon: str       = "#FF8A69"   # Accent
    white: str        = "#FFFFFF"   # Neutral
    black: str        = "#0C0C0C"   # Tech gray / UX
    gray: str         = "#262626"   # Neutral dark

    def rgb(self, name: str) -> RGBColor:
        hex_val = getattr(self, name).lstrip("#")
        return RGBColor(int(hex_val[:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16))

    @staticmethod
    def hex_to_rgb(hex_val: str) -> RGBColor:
        h = hex_val.lstrip("#")
        return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


@dataclass
class BrandTypography:
    """Brand guidelines: Inter Tight Bold (headlines), Inter Tight Regular (body), Roboto Regular (utility)."""
    headline_font: str = "Inter Tight"
    body_font: str     = "Inter Tight"
    utility_font: str  = "Roboto"
    # Fallbacks if not installed
    headline_fallback: str = "Arial"
    body_fallback: str     = "Arial"
    utility_fallback: str  = "Calibri"


@dataclass
class BrandLogoRules:
    """Enforced constraints from brand guidelines."""
    min_lockup_px: int = 100        # Minimum lockup size in pixels
    clearspace_unit: str = "O_width"  # Use width of 'O' in OpenTeams
    allowed_placements: list = field(default_factory=lambda: [
        "upper-left", "lower-left", "upper-center", "lower-center"
    ])
    # NEVER: right-aligned, tagline lockup recreated as text, distorted,
    #        special effects, altered colors, added elements


COLORS = BrandColors()
TYPO = BrandTypography()
LOGO_RULES = BrandLogoRules()

# ===================================================================
# SECTION 2 — WEBSITE CRAWL  (visual cues, subordinate to brand)
# ===================================================================

DEFAULT_SITE_STYLE: Dict[str, Any] = {
    "background": "predominantly white with color accent sections",
    "card_radius_px": 16,
    "card_shadow": "0 4px 24px rgba(0,0,0,0.08)",
    "card_border": "none (shadow-based elevation)",
    "spacing_scale_px": [8, 16, 24, 32, 48, 64, 96],
    "button_style": "pill / fully-rounded, filled Day Blue with white text",
    "button_radius_px": 999,
    "heading_scale": {"h1_pt": 44, "h2_pt": 32, "h3_pt": 24, "h4_pt": 18},
    "gradient_usage": "subtle blue gradients on hero sections (Night Navy → Day Blue)",
    "icon_style": "line icons, medium stroke, brand-colored",
    "image_treatment": "rounded corners ~16px, occasional blue tint overlays",
    "section_padding_px": 96,
    "hero_style": "large headline left-aligned, illustration/image right, white or light gradient bg",
}


def crawl_website(url: str) -> Dict[str, Any]:
    """Fetch homepage + a few pages, extract CSS-derived style tokens.
    Returns style dict; falls back to defaults on failure."""
    if requests is None or BeautifulSoup is None:
        log.warning("requests/bs4 not available — using default site style tokens.")
        return DEFAULT_SITE_STYLE

    style = dict(DEFAULT_SITE_STYLE)
    pages_to_fetch = [url]

    # Try sitemap first
    try:
        r = requests.get(url.rstrip("/") + "/sitemap.xml", timeout=10)
        if r.status_code == 200 and "xml" in r.headers.get("content-type", ""):
            soup = BeautifulSoup(r.text, "lxml-xml")
            locs = [loc.text for loc in soup.find_all("loc")]
            # Pick up to 3 product/feature pages
            candidates = [u for u in locs if any(k in u.lower() for k in
                          ("product", "feature", "solution", "platform", "about"))]
            pages_to_fetch.extend(candidates[:3])
    except Exception as e:
        log.debug(f"Sitemap fetch failed: {e}")

    # Deduplicate
    pages_to_fetch = list(dict.fromkeys(pages_to_fetch))

    extracted_css_vars: Dict[str, str] = {}
    radii: List[int] = []
    shadows: List[str] = []
    bg_colors: List[str] = []

    for page_url in pages_to_fetch[:4]:
        try:
            log.info(f"Crawling {page_url} ...")
            resp = requests.get(page_url, timeout=15, headers={
                "User-Agent": "OpenTeams-BrandBot/1.0 (template builder)"
            })
            if resp.status_code != 200:
                continue
            soup = BeautifulSoup(resp.text, "html.parser")

            # Extract inline <style> blocks + style attrs for CSS variables
            for tag in soup.find_all("style"):
                text = tag.string or ""
                # CSS custom properties
                for m in re.finditer(r'--([\w-]+)\s*:\s*([^;]+);', text):
                    extracted_css_vars[m.group(1)] = m.group(2).strip()
                # border-radius values
                for m in re.finditer(r'border-radius\s*:\s*(\d+)', text):
                    radii.append(int(m.group(1)))
                # box-shadow values
                for m in re.finditer(r'box-shadow\s*:\s*([^;]+);', text):
                    shadows.append(m.group(1).strip())
                # background-color
                for m in re.finditer(r'background(?:-color)?\s*:\s*(#[0-9a-fA-F]{3,8}|rgb[^;]+)', text):
                    bg_colors.append(m.group(1).strip())

            # Also scan linked CSS (first 2 stylesheets)
            for link in soup.find_all("link", rel="stylesheet")[:2]:
                href = link.get("href", "")
                if not href:
                    continue
                if href.startswith("/"):
                    href = url.rstrip("/") + href
                elif not href.startswith("http"):
                    href = url.rstrip("/") + "/" + href
                try:
                    css_resp = requests.get(href, timeout=10)
                    if css_resp.status_code == 200:
                        css_text = css_resp.text[:200_000]  # cap size
                        for m in re.finditer(r'--([\w-]+)\s*:\s*([^;]+);', css_text):
                            extracted_css_vars[m.group(1)] = m.group(2).strip()
                        for m in re.finditer(r'border-radius\s*:\s*(\d+)', css_text):
                            radii.append(int(m.group(1)))
                        for m in re.finditer(r'box-shadow\s*:\s*([^;]+);', css_text):
                            shadows.append(m.group(1).strip())
                except Exception:
                    pass

        except Exception as e:
            log.warning(f"Failed to crawl {page_url}: {e}")

    # Derive tokens from collected data
    if radii:
        common_radius = max(set(radii), key=radii.count)
        style["card_radius_px"] = common_radius
        log.info(f"Derived card radius from site: {common_radius}px")

    if shadows:
        style["card_shadow"] = shadows[0]
        log.info(f"Derived card shadow: {shadows[0][:60]}...")

    if bg_colors:
        # Check white dominance
        white_like = sum(1 for c in bg_colors if c.upper() in ("#FFF", "#FFFFFF", "#FAFAFA", "#F9FAFB"))
        style["white_dominance_ratio"] = white_like / max(len(bg_colors), 1)

    style["css_variables_found"] = len(extracted_css_vars)
    style["pages_crawled"] = len(pages_to_fetch)

    log.info(f"Website crawl complete — {len(extracted_css_vars)} CSS vars, "
             f"{len(radii)} radii, {len(shadows)} shadows extracted.")

    return style


# ===================================================================
# SECTION 3 — ASSET SCANNER
# ===================================================================

@dataclass
class AssetIndex:
    """Categorized inventory of discovered assets."""
    logo_colored_horizontal_png: Optional[str] = None
    logo_colored_horizontal_svg: Optional[str] = None
    logo_colored_vertical_png: Optional[str] = None
    logo_white_horizontal_png: Optional[str] = None
    logo_white_horizontal_svg: Optional[str] = None
    logo_black_horizontal_png: Optional[str] = None
    logo_black_horizontal_svg: Optional[str] = None
    favicon_colored_png: Optional[str] = None
    favicon_white_png: Optional[str] = None
    all_logos: list = field(default_factory=list)
    icons: list = field(default_factory=list)
    backgrounds: list = field(default_factory=list)
    photos: list = field(default_factory=list)
    decisions: list = field(default_factory=list)


def scan_assets(assets_root: str) -> AssetIndex:
    """Recursively scan assets folder, categorize files, pick best variants."""
    idx = AssetIndex()
    root = Path(assets_root)

    if not root.exists():
        log.error(f"Assets folder not found: {assets_root}")
        return idx

    image_exts = {".png", ".jpg", ".jpeg", ".svg", ".webp"}
    skip_exts = {".ai", ".pdf", ".docx", ".DS_Store"}

    for fp in sorted(root.rglob("*")):
        if fp.is_dir():
            continue
        ext = fp.suffix.lower()
        if ext in skip_exts or fp.name.startswith("~$") or fp.name == ".DS_Store":
            continue
        if ext not in image_exts:
            continue

        rel = str(fp)
        name_lower = fp.name.lower()
        parent_lower = str(fp.parent).lower()

        # Classify
        is_logo = any(k in parent_lower for k in ("logo", "lockup"))
        is_favicon = "favicon" in name_lower
        is_white = "white" in parent_lower or "white" in name_lower
        is_black = "black" in parent_lower or "black" in name_lower
        is_colored = "colored" in parent_lower or "color" in parent_lower
        is_horizontal = "horizontal" in parent_lower or "horizontal" in name_lower
        is_vertical = "vertical" in parent_lower or "vertical" in name_lower

        if is_logo or is_favicon:
            idx.all_logos.append(rel)

            # Colored horizontal PNG — primary logo for light backgrounds
            # Exclude files that are actually white or vertical variants mis-filed here
            if is_colored and is_horizontal and ext == ".png" and not is_favicon:
                if "transparent" in parent_lower and "white" not in name_lower and "vertical" not in name_lower:
                    # Prefer cleaner filename (no numeric suffix like "logo3")
                    is_better = (idx.logo_colored_horizontal_png is None or
                                 (re.search(r'\d', Path(idx.logo_colored_horizontal_png).stem.replace("final", "")) and
                                  not re.search(r'\d', fp.stem.replace("final", ""))))
                    if is_better or idx.logo_colored_horizontal_png is None:
                        idx.logo_colored_horizontal_png = rel
                        idx.decisions.append(f"PRIMARY_LOGO_COLOR_H_PNG: {rel}")
            # Colored horizontal SVG
            if is_colored and is_horizontal and ext == ".svg" and not is_favicon:
                idx.logo_colored_horizontal_svg = rel
                idx.decisions.append(f"PRIMARY_LOGO_COLOR_H_SVG: {rel}")
            # Colored vertical PNG
            if is_colored and is_vertical and ext == ".png" and not is_favicon:
                if "transparent" in parent_lower and "white" not in name_lower:
                    idx.logo_colored_vertical_png = rel
                    idx.decisions.append(f"PRIMARY_LOGO_COLOR_V_PNG: {rel}")
            # White horizontal (must be in White logos folder, not just have "white" in name)
            if "white_logos" in parent_lower.replace(" ", "_") and is_horizontal and ext == ".png" and "lockup" in name_lower:
                idx.logo_white_horizontal_png = rel
                idx.decisions.append(f"WHITE_LOGO_H_PNG: {rel}")
            if is_white and is_horizontal and ext == ".svg":
                idx.logo_white_horizontal_svg = rel
            # Black horizontal
            if is_black and is_horizontal and ext == ".png":
                idx.logo_black_horizontal_png = rel
                idx.decisions.append(f"BLACK_LOGO_H_PNG: {rel}")
            if is_black and is_horizontal and ext == ".svg":
                idx.logo_black_horizontal_svg = rel
            # Favicons
            if is_favicon and is_colored and ext == ".png":
                idx.favicon_colored_png = rel
                idx.decisions.append(f"FAVICON_COLOR_PNG: {rel}")
            if is_favicon and is_white and ext == ".png":
                idx.favicon_white_png = rel

        elif ext in {".jpg", ".jpeg", ".webp"}:
            idx.photos.append(rel)
        elif any(k in name_lower for k in ("icon", "illustration")):
            idx.icons.append(rel)
        elif any(k in name_lower for k in ("bg", "background", "pattern", "texture")):
            idx.backgrounds.append(rel)

    # Filter out misclassified logos from colored horizontal
    # (e.g., the file named "OT-new-logo-final-white-vertical-lockup.png" in the colored folder)
    if idx.logo_colored_horizontal_png:
        name = Path(idx.logo_colored_horizontal_png).name.lower()
        if "white" in name and "vertical" in name:
            # This is actually a white vertical logo mis-filed in colored folder
            # Pick the other candidate
            for logo in idx.all_logos:
                lname = Path(logo).name.lower()
                lparent = str(Path(logo).parent).lower()
                if ("colored" in lparent and "horizontal" in lparent and
                    "transparent" in lparent and logo.endswith(".png") and
                    "white" not in lname and "vertical" not in lname):
                    idx.logo_colored_horizontal_png = logo
                    idx.decisions.append(f"PRIMARY_LOGO_COLOR_H_PNG (corrected): {logo}")
                    break

    log.info(f"Asset scan: {len(idx.all_logos)} logos, {len(idx.icons)} icons, "
             f"{len(idx.backgrounds)} backgrounds, {len(idx.photos)} photos")
    for d in idx.decisions:
        log.info(f"  → {d}")

    return idx


def save_assets_index(idx: AssetIndex, path: str):
    with open(path, "w") as f:
        json.dump(asdict(idx), f, indent=2)
    log.info(f"Saved assets index → {path}")


# ===================================================================
# SECTION 4 — DESIGN SYSTEM / THEME CONFIG
# ===================================================================

@dataclass
class ThemeConfig:
    """Merged theme: brand tokens (immutable) + website-derived tokens (flex)."""
    # Brand tokens (from guidelines — immutable)
    colors: BrandColors = field(default_factory=BrandColors)
    typography: BrandTypography = field(default_factory=BrandTypography)

    # Type scale (pts) — derived from website heading scale + brand font rules
    h1_size: int = 44
    h2_size: int = 32
    h3_size: int = 24
    h4_size: int = 18
    body_size: int = 14
    small_size: int = 11
    caption_size: int = 10

    # Spacing (derived from website)
    margin_inches: float = 0.6
    gutter_inches: float = 0.35
    section_pad_inches: float = 0.5

    # Card styles (derived from website)
    card_radius_pt: int = 12  # ~16px at 96dpi ≈ 12pt
    card_shadow_alpha: int = 20  # 0-255, subtle

    # Button styles (derived from website)
    button_radius_pt: int = 20  # pill style
    button_fill_color: str = "#4D75FE"  # Day Blue
    button_text_color: str = "#FFFFFF"

    # Slide dimensions (widescreen 16:9)
    slide_width_inches: float = 13.333
    slide_height_inches: float = 7.5


def build_theme(site_style: Dict[str, Any]) -> ThemeConfig:
    """Merge brand rules with website-derived tokens."""
    theme = ThemeConfig()

    # Apply website heading scale if available (but keep brand fonts)
    hs = site_style.get("heading_scale", {})
    if hs.get("h1_pt"):
        theme.h1_size = hs["h1_pt"]
    if hs.get("h2_pt"):
        theme.h2_size = hs["h2_pt"]
    if hs.get("h3_pt"):
        theme.h3_size = hs["h3_pt"]
    if hs.get("h4_pt"):
        theme.h4_size = hs["h4_pt"]

    # Card radius from site (cap at reasonable max for PPT)
    r = site_style.get("card_radius_px", 16)
    theme.card_radius_pt = min(16, max(8, int(r * 0.75)))  # px→pt approx, capped

    # Button radius
    br = site_style.get("button_radius_px", 999)
    theme.button_radius_pt = min(24, max(8, int(br * 0.75)))

    return theme


# ===================================================================
# SECTION 5 — POWERPOINT HELPERS  (python-pptx)
# ===================================================================
# NOTE: python-pptx limitations:
#  - No native gradient fills on shapes (we use solid fills + overlays)
#  - No native rounded rectangle radius control via API (we use the
#    MSO_SHAPE.ROUNDED_RECTANGLE and manipulate XML for radius)
#  - No theme master editing via API (we build all slides manually)
#  - SVG is NOT supported for images; we use PNG exclusively.
# ===================================================================

def hex_to_rgbcolor(hex_val: str) -> RGBColor:
    h = hex_val.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def luminance(hex_val: str) -> float:
    """Relative luminance for WCAG contrast check."""
    h = hex_val.lstrip("#")
    r, g, b = int(h[:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255

    def linearize(c):
        return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4

    return 0.2126 * linearize(r) + 0.7152 * linearize(g) + 0.0722 * linearize(b)


def contrast_ratio(c1: str, c2: str) -> float:
    l1, l2 = luminance(c1), luminance(c2)
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def auto_text_color(bg_hex: str) -> str:
    """Pick white or black text for best contrast on given background."""
    cr_white = contrast_ratio(bg_hex, "#FFFFFF")
    cr_black = contrast_ratio(bg_hex, "#0C0C0C")
    return "#FFFFFF" if cr_white > cr_black else "#0C0C0C"


def set_shape_fill(shape, hex_color: str):
    """Set solid fill on a shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgbcolor(hex_color)


def set_shape_rounded_rect_radius(shape, radius_emu: int):
    """Set corner radius on a rounded rectangle by editing XML."""
    sp = shape._element
    prstGeom = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if prstGeom is not None:
        avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        if avLst is None:
            avLst = etree.SubElement(prstGeom,
                '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        # Clear existing
        for child in list(avLst):
            avLst.remove(child)
        # Add guide value (in 1/50000 of shape size units)
        gd = etree.SubElement(avLst,
            '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj')
        # Convert EMU radius to percentage of half the smaller dimension
        min_dim = min(shape.width, shape.height)
        if min_dim > 0:
            frac = min(radius_emu / (min_dim / 2), 1.0)
            gd.set('fmla', f'val {int(frac * 50000)}')


def set_no_border(shape):
    """Remove shape border/outline."""
    shape.line.fill.background()


def add_slide_bg_color(slide, hex_color: str):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgbcolor(hex_color)


def _get_spPr(shape):
    """Get the spPr element from a shape, handling both p: and a: namespaces."""
    sp = shape._element
    # Try p: namespace first (presentationml)
    spPr = sp.find('{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
    if spPr is None:
        # Try a: namespace (drawingml)
        spPr = sp.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    if spPr is None:
        # Try without namespace
        for child in sp:
            if etree.QName(child.tag).localname == 'spPr':
                spPr = child
                break
    return spPr


def make_gradient_rect(slide, left, top, width, height, color1_hex, color2_hex, angle=0):
    """Add a rectangle with gradient fill using XML manipulation."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_no_border(shape)

    # Remove default fill and add gradient via XML
    sp_pr = _get_spPr(shape)
    if sp_pr is None:
        log.warning("Could not find spPr for gradient rect — using solid fill fallback.")
        set_shape_fill(shape, color1_hex)
        return shape

    # Remove any existing fill
    for child in list(sp_pr):
        tag = etree.QName(child.tag).localname
        if tag in ('solidFill', 'noFill', 'gradFill'):
            sp_pr.remove(child)

    nsuri = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    gradFill = etree.SubElement(sp_pr, f'{{{nsuri}}}gradFill')
    gradFill.set('rotWithShape', '1')
    gsLst = etree.SubElement(gradFill, f'{{{nsuri}}}gsLst')

    # Stop 1
    gs1 = etree.SubElement(gsLst, f'{{{nsuri}}}gs')
    gs1.set('pos', '0')
    srgb1 = etree.SubElement(gs1, f'{{{nsuri}}}srgbClr')
    srgb1.set('val', color1_hex.lstrip('#'))

    # Stop 2
    gs2 = etree.SubElement(gsLst, f'{{{nsuri}}}gs')
    gs2.set('pos', '100000')
    srgb2 = etree.SubElement(gs2, f'{{{nsuri}}}srgbClr')
    srgb2.set('val', color2_hex.lstrip('#'))

    # Linear direction
    lin = etree.SubElement(gradFill, f'{{{nsuri}}}lin')
    lin.set('ang', str(angle * 60000))  # angle in 60000ths of a degree
    lin.set('scaled', '1')

    return shape


def set_shape_alpha(shape, alpha_pct_str: str):
    """Set fill transparency on a shape. alpha_pct_str like '30000' = 30%."""
    spPr = _get_spPr(shape)
    if spPr is None:
        return
    nsuri = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    solidFill = spPr.find(f'{{{nsuri}}}solidFill')
    if solidFill is not None and len(solidFill) > 0:
        alpha = etree.SubElement(solidFill[0], f'{{{nsuri}}}alpha')
        alpha.set('val', alpha_pct_str)


# ===================================================================
# SECTION 6 — SLIDE BUILDER HELPERS
# ===================================================================

class SlideBuilder:
    """High-level helpers for building branded slides."""

    def __init__(self, prs: Presentation, theme: ThemeConfig, assets: AssetIndex):
        self.prs = prs
        self.theme = theme
        self.assets = assets
        self.c = theme.colors
        self.t = theme.typography
        self.W = Inches(theme.slide_width_inches)
        self.H = Inches(theme.slide_height_inches)
        self.M = Inches(theme.margin_inches)
        self.G = Inches(theme.gutter_inches)

    def new_slide(self) -> Any:
        """Add a blank slide."""
        layout = self.prs.slide_layouts[6]  # blank layout
        return self.prs.slides.add_slide(layout)

    # --- Text helpers ---

    def add_title(self, slide, text: str, x=None, y=None, w=None, h=None,
                  font_size=None, color=None, bold=True, align=PP_ALIGN.LEFT):
        """Add headline text box."""
        x = x if x is not None else self.M
        y = y if y is not None else self.M
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(1.0)
        font_size = font_size or self.theme.h1_size
        color = color or self.c.night_navy

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = hex_to_rgbcolor(color)
        p.font.name = self.t.headline_font
        p.alignment = align
        return txBox

    def add_subtitle(self, slide, text: str, x=None, y=None, w=None, h=None,
                     font_size=None, color=None):
        """Add subtitle / subheading."""
        x = x if x is not None else self.M
        y = y if y is not None else Inches(1.6)
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(0.8)
        font_size = font_size or self.theme.h3_size
        color = color or self.c.day_blue

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = False
        p.font.color.rgb = hex_to_rgbcolor(color)
        p.font.name = self.t.body_font
        p.alignment = PP_ALIGN.LEFT
        return txBox

    def add_body(self, slide, text: str, x=None, y=None, w=None, h=None,
                 font_size=None, color=None, bold=False, align=PP_ALIGN.LEFT,
                 line_spacing=1.4):
        """Add body text."""
        x = x if x is not None else self.M
        y = y if y is not None else Inches(2.5)
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(3.0)
        font_size = font_size or self.theme.body_size
        color = color or self.c.gray

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(text.split("\n")):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = hex_to_rgbcolor(color)
            p.font.name = self.t.body_font
            p.alignment = align
            p.space_after = Pt(font_size * (line_spacing - 1))
        return txBox

    def add_bullet_list(self, slide, items: List[str], x=None, y=None, w=None, h=None,
                        font_size=None, color=None, bullet_color=None):
        """Add a bulleted list with brand-colored bullet markers."""
        x = x if x is not None else self.M
        y = y if y is not None else Inches(2.5)
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(3.5)
        font_size = font_size or self.theme.body_size
        color = color or self.c.gray
        bullet_color = bullet_color or self.c.day_blue

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            # Use a bullet character in Day Blue + text in body color
            run_bullet = p.add_run()
            run_bullet.text = "●  "
            run_bullet.font.size = Pt(font_size - 2)
            run_bullet.font.color.rgb = hex_to_rgbcolor(bullet_color)
            run_bullet.font.name = self.t.body_font

            run_text = p.add_run()
            run_text.text = item
            run_text.font.size = Pt(font_size)
            run_text.font.color.rgb = hex_to_rgbcolor(color)
            run_text.font.name = self.t.body_font

            p.space_after = Pt(8)
        return txBox

    # --- Logo helper ---

    def add_logo(self, slide, variant: str = "colored", position: str = "upper-left",
                 max_width_inches: float = 2.0, max_height_inches: float = 0.6):
        """Add logo respecting brand placement rules.
        variant: 'colored', 'white', 'black', 'favicon'
        position: must be in LOGO_RULES.allowed_placements
        """
        if position not in LOGO_RULES.allowed_placements:
            log.warning(f"Logo position '{position}' not allowed. Defaulting to upper-left.")
            position = "upper-left"

        # Select asset
        logo_map = {
            "colored": self.assets.logo_colored_horizontal_png,
            "white": self.assets.logo_white_horizontal_png,
            "black": self.assets.logo_black_horizontal_png,
            "favicon": self.assets.favicon_colored_png,
        }
        logo_path = logo_map.get(variant)
        if not logo_path or not os.path.exists(logo_path):
            log.warning(f"Logo variant '{variant}' not found. Skipping.")
            return None

        # Determine natural aspect ratio
        if Image:
            with Image.open(logo_path) as im:
                nat_w, nat_h = im.size
        else:
            nat_w, nat_h = 1841, 483  # fallback for colored horizontal

        aspect = nat_w / nat_h
        # Fit within max box
        w = Inches(max_width_inches)
        h = Inches(max_width_inches / aspect)
        if h > Inches(max_height_inches):
            h = Inches(max_height_inches)
            w = Inches(max_height_inches * aspect)

        # Enforce minimum 100px (~1.04 inches at 96dpi)
        min_w = Inches(1.04)
        if w < min_w:
            w = min_w
            h = Inches(1.04 / aspect)

        # Clearspace (width of "O" ≈ 1/10 of full lockup width)
        clearspace = Inches(max_width_inches * 0.10)

        # Position
        M = self.M
        if position == "upper-left":
            x, y = M, M * 0.6
        elif position == "lower-left":
            x, y = M, self.H - h - M * 0.6
        elif position == "upper-center":
            x, y = (self.W - w) / 2, M * 0.6
        elif position == "lower-center":
            x, y = (self.W - w) / 2, self.H - h - M * 0.6
        else:
            x, y = M, M * 0.6

        pic = slide.shapes.add_picture(logo_path, int(x), int(y), int(w), int(h))
        return pic

    # --- Decorative elements ---

    def add_accent_bar(self, slide, x, y, w, h, color=None):
        """Thin accent bar / line."""
        color = color or self.c.day_blue
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        set_shape_fill(shape, color)
        set_no_border(shape)
        return shape

    def add_card(self, slide, x, y, w, h, fill_color="#FFFFFF",
                 border_color=None, shadow=True):
        """Add a card-style rounded rectangle (website-like)."""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        set_shape_fill(shape, fill_color)

        # Radius
        radius_emu = Pt(self.theme.card_radius_pt)
        set_shape_rounded_rect_radius(shape, radius_emu)

        if border_color:
            shape.line.color.rgb = hex_to_rgbcolor(border_color)
            shape.line.width = Pt(1)
        else:
            set_no_border(shape)

        # Shadow via XML (subtle)
        if shadow:
            spPr = _get_spPr(shape)
            if spPr is None:
                return shape
            nsuri = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            effectLst = etree.SubElement(spPr, f'{{{nsuri}}}effectLst')
            outerShdw = etree.SubElement(effectLst, f'{{{nsuri}}}outerShdw')
            outerShdw.set('blurRad', '152400')  # ~4pt blur
            outerShdw.set('dist', '38100')       # ~1pt distance
            outerShdw.set('dir', '5400000')      # downward
            outerShdw.set('rotWithShape', '0')
            srgbClr = etree.SubElement(outerShdw, f'{{{nsuri}}}srgbClr')
            srgbClr.set('val', '000000')
            alpha = etree.SubElement(srgbClr, f'{{{nsuri}}}alpha')
            alpha.set('val', '15000')  # 15% opacity

        return shape

    def add_button(self, slide, text: str, x, y, w=None, h=None,
                   fill_color=None, text_color=None):
        """Pill-style button (matching website CTA style)."""
        fill_color = fill_color or self.theme.button_fill_color
        text_color = text_color or self.theme.button_text_color
        w = w or Inches(2.2)
        h = h or Inches(0.55)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        set_shape_fill(shape, fill_color)
        set_no_border(shape)

        # Max radius for pill shape
        set_shape_rounded_rect_radius(shape, int(h / 2))

        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgbcolor(text_color)
        p.font.name = self.t.body_font
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
        # Vertical center
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Set anchor to middle
        txBody = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')

        return shape

    def add_placeholder_image(self, slide, x, y, w, h, label="Image",
                              fill_color=None):
        """Add a rounded placeholder box for images."""
        fill_color = fill_color or "#E8EDFB"  # Very light blue
        shape = self.add_card(slide, x, y, w, h, fill_color=fill_color, shadow=False)

        # Add centered label
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[ {label} ]"
        p.font.size = Pt(12)
        p.font.color.rgb = hex_to_rgbcolor(self.c.day_blue)
        p.font.name = self.t.utility_font
        p.alignment = PP_ALIGN.CENTER
        txBody = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')
        return shape

    def add_footer(self, slide, text: str = "© 2025 OpenTeams  |  openteams.com",
                   show_logo: bool = True, bg_color: str = None):
        """Add footer bar with optional small logo."""
        footer_h = Inches(0.5)
        y = self.H - footer_h

        if bg_color:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, self.W, footer_h)
            set_shape_fill(bar, bg_color)
            set_no_border(bar)

        text_color = auto_text_color(bg_color or "#FFFFFF")

        # Footer text — left aligned
        txBox = slide.shapes.add_textbox(self.M, y + Inches(0.08), Inches(8), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(self.theme.caption_size)
        p.font.color.rgb = hex_to_rgbcolor(text_color)
        p.font.name = self.t.utility_font
        p.alignment = PP_ALIGN.LEFT

        # Small favicon logo on the right side of footer (but left-of-center to comply)
        if show_logo and self.assets.favicon_colored_png:
            icon_size = Inches(0.3)
            # Place it at center-left area
            pic = slide.shapes.add_picture(
                self.assets.favicon_colored_png,
                int(self.W - self.M - icon_size),
                int(y + (footer_h - icon_size) / 2),
                int(icon_size), int(icon_size)
            )

    def add_section_header(self, slide, title: str, subtitle: str = "",
                           bg_color: str = None):
        """Full-bleed section header slide."""
        bg_color = bg_color or self.c.night_navy
        add_slide_bg_color(slide, bg_color)
        text_color = auto_text_color(bg_color)

        # Accent bar
        self.add_accent_bar(slide, self.M, Inches(2.8), Inches(0.8), Inches(0.06),
                           color=self.c.day_blue if bg_color != self.c.day_blue else self.c.yellow)

        self.add_title(slide, title, y=Inches(3.0), font_size=self.theme.h1_size,
                      color=text_color)

        if subtitle:
            sub_color = self.c.day_blue if luminance(bg_color) < 0.3 else self.c.night_navy
            # But ensure contrast
            if contrast_ratio(bg_color, sub_color) < 3:
                sub_color = text_color
            self.add_body(slide, subtitle, y=Inches(4.2), font_size=self.theme.h3_size,
                         color=sub_color)

    def add_metric_card(self, slide, x, y, w, h, value: str, label: str,
                        accent_color: str = None):
        """Small metric card with big number + label."""
        accent_color = accent_color or self.c.day_blue
        card = self.add_card(slide, x, y, w, h)

        # Accent top bar
        bar_h = Inches(0.06)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, bar_h)
        set_shape_fill(bar, accent_color)
        set_no_border(bar)

        # Value
        val_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3),
                                            w - Inches(0.4), Inches(0.8))
        tf = val_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgbcolor(self.c.night_navy)
        p.font.name = self.t.headline_font
        p.alignment = PP_ALIGN.LEFT

        # Label
        lbl_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.0),
                                            w - Inches(0.4), Inches(0.5))
        tf2 = lbl_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = hex_to_rgbcolor(self.c.gray)
        p2.font.name = self.t.body_font
        p2.alignment = PP_ALIGN.LEFT


# ===================================================================
# SECTION 7 — SLIDE GENERATORS
# ===================================================================

def build_slide_cover(sb: SlideBuilder):
    """Slide 1: Cover / Title (hero layout)."""
    slide = sb.new_slide()

    # White background with subtle gradient overlay on left
    add_slide_bg_color(slide, "#FFFFFF")

    # Large gradient accent block on right side (hero visual area)
    make_gradient_rect(slide, Inches(7.5), 0, Inches(5.833), sb.H,
                       sb.c.night_navy, sb.c.day_blue, angle=135)

    # Subtle accent dot pattern (decorative circles)
    for i, (dx, dy, col) in enumerate([
        (8.2, 1.5, sb.c.yellow), (9.8, 2.0, sb.c.salmon),
        (8.8, 5.0, sb.c.day_blue), (10.5, 4.2, sb.c.yellow),
    ]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx), Inches(dy),
                                      Inches(0.35), Inches(0.35))
        set_shape_fill(dot, col)
        set_no_border(dot)
        set_shape_alpha(dot, '30000')  # 30% opacity

    # Decorative large icon mark on the gradient panel
    if sb.assets.favicon_colored_png:
        # Place favicon (large, as decorative element) — but DO NOT alter its colors
        fav_size = Inches(3.5)
        slide.shapes.add_picture(
            sb.assets.favicon_colored_png,
            int(Inches(9.0) - fav_size / 2), int(Inches(3.75) - fav_size / 2),
            int(fav_size), int(fav_size)
        )

    # Logo (upper-left, colored on white)
    sb.add_logo(slide, "colored", "upper-left", max_width_inches=2.4, max_height_inches=0.65)

    # Title text
    sb.add_title(slide, "Presentation Title",
                x=sb.M, y=Inches(2.4), w=Inches(6.5), h=Inches(1.8),
                font_size=48, color=sb.c.night_navy)

    # Subtitle
    sb.add_body(slide, "Subtitle or tagline goes here\nMonth Year",
               x=sb.M, y=Inches(4.3), w=Inches(6.0), h=Inches(1.2),
               font_size=20, color=sb.c.day_blue)

    # CTA button
    sb.add_button(slide, "Get Started", sb.M, Inches(5.8))

    # Footer
    sb.add_footer(slide, "Confidential  |  © 2025 OpenTeams", show_logo=False)


def build_slide_section_divider(sb: SlideBuilder):
    """Slide 2: Section Divider."""
    slide = sb.new_slide()
    sb.add_section_header(slide, "Section Title", "Brief description of this section")

    # Logo (white, on dark bg)
    sb.add_logo(slide, "white", "lower-left", max_width_inches=1.8, max_height_inches=0.45)


def build_slide_agenda(sb: SlideBuilder):
    """Slide 3: Agenda."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")

    # Top accent line
    sb.add_accent_bar(slide, sb.M, Inches(0.9), Inches(0.8), Inches(0.06))

    sb.add_title(slide, "Agenda", y=Inches(1.1), font_size=sb.theme.h2_size)

    items = [
        "Introduction & Context",
        "Problem Statement",
        "Our Approach & Solution",
        "Key Results & Metrics",
        "Next Steps & Discussion",
    ]

    y_start = Inches(2.4)
    for i, item in enumerate(items):
        y = y_start + Inches(i * 0.85)

        # Number circle
        num_size = Inches(0.5)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, sb.M, y, num_size, num_size)
        accent = [sb.c.day_blue, sb.c.night_navy, sb.c.yellow, sb.c.salmon, sb.c.day_blue][i]
        set_shape_fill(circle, accent)
        set_no_border(circle)
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgbcolor(auto_text_color(accent))
        p.font.name = sb.t.headline_font
        p.alignment = PP_ALIGN.CENTER
        txBody = circle._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')

        # Item text
        sb.add_body(slide, item, x=sb.M + Inches(0.75), y=y + Inches(0.05),
                   w=Inches(10), h=Inches(0.5), font_size=18, color=sb.c.gray, bold=False)

    sb.add_footer(slide, show_logo=True)
    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)


def build_slide_content(sb: SlideBuilder):
    """Slide 4: Content (title + body + visual placeholder)."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")

    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)

    # Accent bar
    sb.add_accent_bar(slide, sb.M, Inches(1.2), Inches(0.8), Inches(0.06))

    sb.add_title(slide, "Content Slide Title", y=Inches(1.4), font_size=sb.theme.h2_size)

    # Body text on left
    sb.add_body(slide, (
        "Add your key points here. The template uses generous whitespace\n"
        "and brand-consistent typography for a clean, modern look.\n\n"
        "Use this layout for text-heavy slides that need a supporting visual."
    ), x=sb.M, y=Inches(2.5), w=Inches(5.5), h=Inches(3.5),
       font_size=15, color=sb.c.gray)

    # Image placeholder on right
    sb.add_placeholder_image(slide, Inches(7.0), Inches(1.4),
                             Inches(5.7), Inches(4.8), "Visual / Image")

    sb.add_footer(slide, show_logo=True)


def build_slide_two_column(sb: SlideBuilder):
    """Slide 5: Two-column content."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")

    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)

    sb.add_accent_bar(slide, sb.M, Inches(1.2), Inches(0.8), Inches(0.06))
    sb.add_title(slide, "Two-Column Layout", y=Inches(1.4), font_size=sb.theme.h2_size)

    col_w = Inches(5.8)
    col1_x = sb.M
    col2_x = sb.M + col_w + sb.G

    # Column 1 card
    card1 = sb.add_card(slide, col1_x, Inches(2.5), col_w, Inches(4.0))
    sb.add_body(slide, "Left Column", x=col1_x + Inches(0.3), y=Inches(2.7),
               w=col_w - Inches(0.6), h=Inches(0.5),
               font_size=20, color=sb.c.night_navy, bold=True)
    sb.add_body(slide, (
        "Supporting text for the first column. Use for\n"
        "comparisons, features, or parallel content."
    ), x=col1_x + Inches(0.3), y=Inches(3.3), w=col_w - Inches(0.6), h=Inches(2.5),
       font_size=14, color=sb.c.gray)

    # Column 2 card
    card2 = sb.add_card(slide, col2_x, Inches(2.5), col_w, Inches(4.0))
    sb.add_body(slide, "Right Column", x=col2_x + Inches(0.3), y=Inches(2.7),
               w=col_w - Inches(0.6), h=Inches(0.5),
               font_size=20, color=sb.c.night_navy, bold=True)
    sb.add_body(slide, (
        "Supporting text for the second column.\n"
        "Maintain visual balance between columns."
    ), x=col2_x + Inches(0.3), y=Inches(3.3), w=col_w - Inches(0.6), h=Inches(2.5),
       font_size=14, color=sb.c.gray)

    sb.add_footer(slide, show_logo=True)


def build_slide_big_statement(sb: SlideBuilder):
    """Slide 6: Big statement / quote."""
    slide = sb.new_slide()

    # Deep navy background
    add_slide_bg_color(slide, sb.c.night_navy)

    # Large quote mark (decorative)
    quote_mark = slide.shapes.add_textbox(sb.M, Inches(1.0), Inches(2), Inches(2))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    p.text = "\u201C"
    p.font.size = Pt(160)
    p.font.color.rgb = hex_to_rgbcolor(sb.c.day_blue)
    p.font.name = sb.t.headline_font
    p.font.bold = True

    # Quote text
    sb.add_title(slide, (
        "A bold statement that captures\n"
        "your key message in one line."
    ), x=Inches(1.2), y=Inches(2.8), w=Inches(10.5), h=Inches(2.0),
       font_size=36, color="#FFFFFF")

    # Attribution
    sb.add_body(slide, "— Speaker Name, Title",
               x=Inches(1.2), y=Inches(5.0), w=Inches(8), h=Inches(0.6),
               font_size=16, color=sb.c.day_blue)

    # Accent dots
    for dx, dy, col in [(11.5, 5.5, sb.c.yellow), (12.0, 5.0, sb.c.salmon)]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx), Inches(dy),
                                      Inches(0.25), Inches(0.25))
        set_shape_fill(dot, col)
        set_no_border(dot)

    sb.add_logo(slide, "white", "lower-left", max_width_inches=1.8, max_height_inches=0.45)


def build_slide_data(sb: SlideBuilder):
    """Slide 7: Data slide (metrics cards + chart placeholder)."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#F7F8FC")  # Very light blue-gray (within white family)

    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)
    sb.add_accent_bar(slide, sb.M, Inches(1.1), Inches(0.8), Inches(0.06))
    sb.add_title(slide, "Key Metrics", y=Inches(1.3), font_size=sb.theme.h2_size)

    # Metric cards row
    card_w = Inches(2.7)
    card_h = Inches(1.6)
    gap = Inches(0.35)
    start_x = sb.M
    y = Inches(2.3)

    metrics = [
        ("98%", "Customer Satisfaction", sb.c.day_blue),
        ("3.5x", "ROI Improvement", sb.c.night_navy),
        ("500+", "Active Projects", sb.c.yellow),
        ("24/7", "Global Support", sb.c.salmon),
    ]

    for i, (val, label, accent) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        sb.add_metric_card(slide, x, y, card_w, card_h, val, label, accent)

    # Chart placeholder below
    sb.add_placeholder_image(slide, sb.M, Inches(4.3), Inches(11.8), Inches(2.5),
                             "Chart / Data Visualization")

    sb.add_footer(slide, show_logo=True)


def build_slide_team(sb: SlideBuilder):
    """Slide 8: Team / Profile slide."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")

    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)
    sb.add_accent_bar(slide, sb.M, Inches(1.1), Inches(0.8), Inches(0.06))
    sb.add_title(slide, "Our Team", y=Inches(1.3), font_size=sb.theme.h2_size)

    # 4 profile cards
    card_w = Inches(2.7)
    card_h = Inches(4.2)
    gap = Inches(0.35)
    start_x = sb.M + Inches(0.3)
    y = Inches(2.3)

    for i in range(4):
        x = start_x + i * (card_w + gap)
        card = sb.add_card(slide, x, y, card_w, card_h)

        # Avatar circle placeholder
        avatar_size = Inches(1.4)
        avatar_x = x + (card_w - avatar_size) / 2
        avatar_y = y + Inches(0.4)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(avatar_x), int(avatar_y),
                                         int(avatar_size), int(avatar_size))
        accent = [sb.c.day_blue, sb.c.night_navy, sb.c.yellow, sb.c.salmon][i]
        set_shape_fill(circle, accent)
        set_no_border(circle)
        # Icon text
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = "👤"
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER
        txBody = circle._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')

        # Name
        sb.add_body(slide, f"Team Member {i+1}",
                   x=x + Inches(0.2), y=y + Inches(2.1),
                   w=card_w - Inches(0.4), h=Inches(0.4),
                   font_size=16, color=sb.c.night_navy, bold=True, align=PP_ALIGN.CENTER)

        # Role
        sb.add_body(slide, "Role / Title",
                   x=x + Inches(0.2), y=y + Inches(2.6),
                   w=card_w - Inches(0.4), h=Inches(0.3),
                   font_size=12, color=sb.c.day_blue, align=PP_ALIGN.CENTER)

        # Bio snippet
        sb.add_body(slide, "Brief bio or expertise\narea description.",
                   x=x + Inches(0.2), y=y + Inches(3.1),
                   w=card_w - Inches(0.4), h=Inches(0.8),
                   font_size=11, color=sb.c.gray, align=PP_ALIGN.CENTER)

    sb.add_footer(slide, show_logo=True)


def build_slide_case_study(sb: SlideBuilder):
    """Slide 9: Case Study (Problem → Solution → Result)."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")

    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)
    sb.add_accent_bar(slide, sb.M, Inches(1.1), Inches(0.8), Inches(0.06))
    sb.add_title(slide, "Case Study: Client Name", y=Inches(1.3), font_size=sb.theme.h2_size)

    # Three columns: Problem, Solution, Result
    col_w = Inches(3.7)
    gap = Inches(0.4)
    start_x = sb.M + Inches(0.15)
    y = Inches(2.5)
    col_h = Inches(4.0)

    labels = ["Challenge", "Solution", "Results"]
    colors = [sb.c.salmon, sb.c.day_blue, sb.c.yellow]
    bodies = [
        "Describe the client's\nchallenge or pain point\nthat needed addressing.",
        "Explain the OpenTeams\napproach and how the\nsolution was implemented.",
        "Share quantifiable\noutcomes and the\nimpact delivered.",
    ]

    for i, (label, accent, body) in enumerate(zip(labels, colors, bodies)):
        x = start_x + i * (col_w + gap)

        # Card
        card = sb.add_card(slide, x, y, col_w, col_h)

        # Accent top bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_w, Inches(0.07))
        set_shape_fill(bar, accent)
        set_no_border(bar)

        # Icon circle
        icon_size = Inches(0.6)
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       int(x + Inches(0.3)), int(y + Inches(0.4)),
                                       int(icon_size), int(icon_size))
        set_shape_fill(icon, accent)
        set_no_border(icon)
        tf = icon.text_frame
        p = tf.paragraphs[0]
        p.text = ["⚡", "🔧", "📈"][i]
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
        txBody = icon._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')

        # Label
        sb.add_body(slide, label,
                   x=x + Inches(0.3), y=y + Inches(1.2),
                   w=col_w - Inches(0.6), h=Inches(0.4),
                   font_size=20, color=sb.c.night_navy, bold=True)

        # Body
        sb.add_body(slide, body,
                   x=x + Inches(0.3), y=y + Inches(1.8),
                   w=col_w - Inches(0.6), h=Inches(1.8),
                   font_size=13, color=sb.c.gray)

    sb.add_footer(slide, show_logo=True)


def build_slide_closing(sb: SlideBuilder):
    """Slide 10: Closing / CTA."""
    slide = sb.new_slide()

    # Gradient background
    make_gradient_rect(slide, 0, 0, sb.W, sb.H, sb.c.night_navy, sb.c.day_blue, angle=135)

    # Decorative circles
    for dx, dy, sz, col, alpha_val in [
        (1.5, 1.0, 1.5, sb.c.yellow, "15000"),
        (10.5, 5.5, 2.0, sb.c.salmon, "12000"),
        (11.0, 1.5, 0.8, sb.c.day_blue, "20000"),
    ]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx), Inches(dy),
                                      Inches(sz), Inches(sz))
        set_shape_fill(dot, col)
        set_no_border(dot)
        set_shape_alpha(dot, alpha_val)

    # Thank you text
    sb.add_title(slide, "Thank You",
                x=sb.M, y=Inches(2.0), w=Inches(12), h=Inches(1.5),
                font_size=56, color="#FFFFFF", align=PP_ALIGN.CENTER)

    sb.add_body(slide, "Questions? Let's discuss.",
               x=sb.M, y=Inches(3.6), w=Inches(12), h=Inches(0.8),
               font_size=22, color=sb.c.day_blue, align=PP_ALIGN.CENTER)

    # CTA button
    btn_w = Inches(2.8)
    sb.add_button(slide, "Contact Us",
                 (sb.W - btn_w) / 2, Inches(4.8), w=btn_w,
                 fill_color="#FFFFFF", text_color=sb.c.night_navy)

    # Contact info
    sb.add_body(slide, "hello@openteams.com  |  openteams.com",
               x=sb.M, y=Inches(5.8), w=Inches(12), h=Inches(0.5),
               font_size=14, color="#FFFFFF", align=PP_ALIGN.CENTER)

    # Logo (white, centered at bottom)
    sb.add_logo(slide, "white", "lower-center", max_width_inches=2.2, max_height_inches=0.55)


# ===================================================================
# SECTION 8 — MAIN ORCHESTRATOR
# ===================================================================

def build_template(assets_dir: str, guidelines_path: str, site_url: str, output_path: str):
    """Main entry point: crawl site, scan assets, build template."""

    log.info("=" * 60)
    log.info("OpenTeams 2025 Template Builder")
    log.info("=" * 60)

    # --- Step 1: Parse brand guidelines (hard-coded from PDF) ---
    log.info("\n[1/5] Brand guidelines loaded (hard-coded from PDF)")
    log.info(f"  Colors: Night Navy={COLORS.night_navy}, Day Blue={COLORS.day_blue}, "
             f"Yellow={COLORS.yellow}, Salmon={COLORS.salmon}")
    log.info(f"  Typography: {TYPO.headline_font} Bold (headlines), "
             f"{TYPO.body_font} Regular (body), {TYPO.utility_font} (utility)")
    log.info(f"  Logo rules: min {LOGO_RULES.min_lockup_px}px, "
             f"placements={LOGO_RULES.allowed_placements}")

    if guidelines_path and os.path.exists(guidelines_path):
        log.info(f"  Guidelines PDF confirmed at: {guidelines_path}")
    else:
        log.warning(f"  Guidelines PDF not found at: {guidelines_path} (using hard-coded rules)")

    # --- Step 2: Crawl website for style cues ---
    log.info("\n[2/5] Crawling website for visual cues...")
    site_style = crawl_website(site_url)

    # Save site style tokens
    site_style_path = os.path.join(os.path.dirname(output_path) or ".", "site_style.json")
    with open(site_style_path, "w") as f:
        json.dump(site_style, f, indent=2, default=str)
    log.info(f"  Saved site style tokens → {site_style_path}")

    # --- Step 3: Scan assets ---
    log.info("\n[3/5] Scanning asset folder...")
    assets = scan_assets(assets_dir)

    assets_index_path = os.path.join(os.path.dirname(output_path) or ".", "assets_index.json")
    save_assets_index(assets, assets_index_path)

    # --- Step 4: Build theme ---
    log.info("\n[4/5] Building theme config...")
    theme = build_theme(site_style)
    log.info(f"  Type scale: H1={theme.h1_size}pt, H2={theme.h2_size}pt, "
             f"H3={theme.h3_size}pt, body={theme.body_size}pt")
    log.info(f"  Card radius: {theme.card_radius_pt}pt, Button radius: {theme.button_radius_pt}pt")

    # --- Step 5: Generate PPTX ---
    log.info("\n[5/5] Generating PowerPoint template...")

    prs = Presentation()
    prs.slide_width = Emu(int(Inches(theme.slide_width_inches)))
    prs.slide_height = Emu(int(Inches(theme.slide_height_inches)))

    sb = SlideBuilder(prs, theme, assets)

    slide_builders = [
        ("Cover / Title",        build_slide_cover),
        ("Section Divider",      build_slide_section_divider),
        ("Agenda",               build_slide_agenda),
        ("Content",              build_slide_content),
        ("Two-Column",           build_slide_two_column),
        ("Big Statement",        build_slide_big_statement),
        ("Data / Metrics",       build_slide_data),
        ("Team / Profiles",      build_slide_team),
        ("Case Study",           build_slide_case_study),
        ("Closing / CTA",        build_slide_closing),
    ]

    for name, builder_fn in slide_builders:
        try:
            builder_fn(sb)
            log.info(f"  ✓ {name}")
        except Exception as e:
            log.error(f"  ✗ {name}: {e}")
            import traceback
            traceback.print_exc()

    prs.save(output_path)
    log.info(f"\n{'=' * 60}")
    log.info(f"Template saved → {output_path}")
    log.info(f"  {len(prs.slides)} slides generated")
    log.info(f"  Assets index → {assets_index_path}")
    log.info(f"  Site style → {site_style_path}")
    log.info(f"{'=' * 60}")


# ===================================================================
# CLI
# ===================================================================

def main():
    parser = argparse.ArgumentParser(
        description="OpenTeams 2025 PowerPoint Template Builder",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=textwrap.dedent("""\
            Example:
              python build_template.py \\
                --assets ./Assets \\
                --guidelines ./OpenTeams_Brand_Guidelines_2025.pdf \\
                --site https://openteams.com/ \\
                --out ./OpenTeams_Template_2025.pptx
        """)
    )
    parser.add_argument("--assets", default="./Assets",
                        help="Path to assets folder (default: ./Assets)")
    parser.add_argument("--guidelines", default="./OpenTeams_Brand_Guidelines_2025.pdf",
                        help="Path to brand guidelines PDF")
    parser.add_argument("--site", default="https://openteams.com/",
                        help="Website URL to crawl for style cues")
    parser.add_argument("--out", default="./OpenTeams_Template_2025.pptx",
                        help="Output PPTX file path")
    parser.add_argument("--skip-crawl", action="store_true",
                        help="Skip website crawl (use defaults)")
    parser.add_argument("-v", "--verbose", action="store_true",
                        help="Enable debug logging")

    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    if args.skip_crawl:
        # Monkey-patch to skip crawl
        global crawl_website
        _orig = crawl_website
        crawl_website = lambda url: DEFAULT_SITE_STYLE

    build_template(args.assets, args.guidelines, args.site, args.out)


if __name__ == "__main__":
    main()

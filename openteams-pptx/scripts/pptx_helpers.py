"""
Low-level PowerPoint shape/gradient/shadow helpers.
Extracted from build_template.py Section 5.
"""
from __future__ import annotations

import logging
import re

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

log = logging.getLogger(__name__)

_HEX_RE = re.compile(r'^#?([0-9a-fA-F]{6})$')
_FALLBACK_COLOR = "#000000"


# ---------------------------------------------------------------------------
# Color utilities
# ---------------------------------------------------------------------------

def hex_to_rgbcolor(hex_val: str) -> RGBColor:
    """Convert a hex color string to an RGBColor.

    Validates input and falls back to black with a warning on bad values.
    """
    m = _HEX_RE.match(hex_val.strip())
    if not m:
        log.warning("Invalid hex color '%s' — falling back to %s", hex_val, _FALLBACK_COLOR)
        hex_val = _FALLBACK_COLOR
    h = hex_val.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def luminance(hex_val: str) -> float:
    """Relative luminance for WCAG contrast check."""
    m = _HEX_RE.match(hex_val.strip())
    if not m:
        return 0.0
    h = hex_val.lstrip("#")
    r, g, b = int(h[:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255

    def linearize(c):
        return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4

    return 0.2126 * linearize(r) + 0.7152 * linearize(g) + 0.0722 * linearize(b)


def contrast_ratio(c1: str, c2: str) -> float:
    l1, l2 = luminance(c1), luminance(c2)
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def auto_text_color(bg_hex: str) -> str:
    """Pick white or black text for best contrast on given background."""
    cr_white = contrast_ratio(bg_hex, "#FFFFFF")
    cr_black = contrast_ratio(bg_hex, "#000000")
    return "#FFFFFF" if cr_white > cr_black else "#000000"


# ---------------------------------------------------------------------------
# Shape fill / border helpers
# ---------------------------------------------------------------------------

def set_shape_fill(shape, hex_color: str):
    """Set solid fill on a shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgbcolor(hex_color)


def set_shape_rounded_rect_radius(shape, radius_emu: int):
    """Set corner radius on a rounded rectangle by editing XML."""
    sp = shape._element
    prstGeom = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if prstGeom is not None:
        avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        if avLst is None:
            avLst = etree.SubElement(prstGeom,
                '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        for child in list(avLst):
            avLst.remove(child)
        gd = etree.SubElement(avLst,
            '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj')
        min_dim = min(shape.width, shape.height)
        if min_dim > 0:
            frac = min(radius_emu / (min_dim / 2), 0.98)  # cap below 1.0 to avoid fully rounded
            gd.set('fmla', f'val {int(frac * 50000)}')


def set_no_border(shape):
    """Remove shape border/outline."""
    shape.line.fill.background()


def add_slide_bg_color(slide, hex_color: str):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgbcolor(hex_color)


def get_spPr(shape):
    """Get the spPr element from a shape, handling both p: and a: namespaces."""
    sp = shape._element
    spPr = sp.find('{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
    if spPr is None:
        spPr = sp.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    if spPr is None:
        for child in sp:
            if etree.QName(child.tag).localname == 'spPr':
                spPr = child
                break
    return spPr


# Keep old name as alias for backward compatibility
_get_spPr = get_spPr


def make_gradient_rect(slide, left, top, width, height, color1_hex, color2_hex, angle=0):
    """Add a rectangle with gradient fill using XML manipulation."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_no_border(shape)

    sp_pr = get_spPr(shape)
    if sp_pr is None:
        set_shape_fill(shape, color1_hex)
        return shape

    for child in list(sp_pr):
        tag = etree.QName(child.tag).localname
        if tag in ('solidFill', 'noFill', 'gradFill'):
            sp_pr.remove(child)

    nsuri = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    gradFill = etree.SubElement(sp_pr, f'{{{nsuri}}}gradFill')
    gradFill.set('rotWithShape', '1')
    gsLst = etree.SubElement(gradFill, f'{{{nsuri}}}gsLst')

    gs1 = etree.SubElement(gsLst, f'{{{nsuri}}}gs')
    gs1.set('pos', '0')
    srgb1 = etree.SubElement(gs1, f'{{{nsuri}}}srgbClr')
    srgb1.set('val', color1_hex.lstrip('#'))

    gs2 = etree.SubElement(gsLst, f'{{{nsuri}}}gs')
    gs2.set('pos', '100000')
    srgb2 = etree.SubElement(gs2, f'{{{nsuri}}}srgbClr')
    srgb2.set('val', color2_hex.lstrip('#'))

    lin = etree.SubElement(gradFill, f'{{{nsuri}}}lin')
    lin.set('ang', str(angle * 60000))
    lin.set('scaled', '1')

    return shape


def set_shape_alpha(shape, alpha_pct_str: str):
    """Set fill transparency on a shape. alpha_pct_str like '30000' = 30%.

    Works with both solid fills and gradient fills.
    """
    spPr = get_spPr(shape)
    if spPr is None:
        return
    nsuri = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Try solid fill first
    solidFill = spPr.find(f'{{{nsuri}}}solidFill')
    if solidFill is not None and len(solidFill) > 0:
        alpha = etree.SubElement(solidFill[0], f'{{{nsuri}}}alpha')
        alpha.set('val', alpha_pct_str)
        return

    # Fall back to gradient fill stops
    gradFill = spPr.find(f'{{{nsuri}}}gradFill')
    if gradFill is not None:
        for gs in gradFill.findall(f'.//{{{nsuri}}}gs'):
            color_elem = gs.find(f'{{{nsuri}}}srgbClr')
            if color_elem is not None:
                alpha = etree.SubElement(color_elem, f'{{{nsuri}}}alpha')
                alpha.set('val', alpha_pct_str)

"""
SlideBuilder — high-level helpers for building branded slides.
Extracted from slide_renderers.py for maintainability.
"""
from __future__ import annotations

import datetime
import os
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

from pptx_helpers import (
    hex_to_rgbcolor, luminance, contrast_ratio, auto_text_color,
    set_shape_fill, set_shape_rounded_rect_radius, set_no_border,
    add_slide_bg_color, make_gradient_rect, set_shape_alpha, get_spPr,
)
from brand_engine import ThemeConfig

_CURRENT_YEAR = datetime.date.today().year

try:
    from PIL import Image
except ImportError:
    Image = None


class SlideBuilder:
    """High-level helpers for building branded slides."""

    # Rotating accent colors used by agenda, metrics, team, case_study renderers
    ACCENT_ROTATION = None  # Set after __init__ from theme colors

    def __init__(self, prs: Presentation, theme: ThemeConfig):
        self.prs = prs
        self.theme = theme
        self.c = theme  # color shortcuts via theme
        self.t = theme  # typography shortcuts via theme
        self.W = Inches(theme.slide_width_inches)
        self.H = Inches(theme.slide_height_inches)
        self.M = Inches(theme.margin_inches)
        self.G = Inches(theme.gutter_inches)
        self.ACCENT_ROTATION = [
            theme.day_blue, theme.night_navy, theme.yellow, theme.salmon,
        ]

    def new_slide(self):
        """Add a blank slide."""
        layout = self.prs.slide_layouts[6]  # blank layout
        return self.prs.slides.add_slide(layout)

    # --- Text helpers ---

    def add_title(self, slide, text: str, x=None, y=None, w=None, h=None,
                  font_size=None, color=None, bold=True, align=PP_ALIGN.LEFT):
        x = x if x is not None else self.M
        y = y if y is not None else self.M
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(1.0)
        font_size = font_size or self.theme.h1_size
        color = color or self.theme.night_navy

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = hex_to_rgbcolor(color)
        p.font.name = self.theme.headline_font
        p.alignment = align
        return txBox

    def add_subtitle(self, slide, text: str, x=None, y=None, w=None, h=None,
                     font_size=None, color=None):
        x = x if x is not None else self.M
        y = y if y is not None else Inches(1.6)
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(0.8)
        font_size = font_size or self.theme.h3_size
        color = color or self.theme.day_blue

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = False
        p.font.color.rgb = hex_to_rgbcolor(color)
        p.font.name = self.theme.body_font
        p.alignment = PP_ALIGN.LEFT
        return txBox

    def add_body(self, slide, text: str, x=None, y=None, w=None, h=None,
                 font_size=None, color=None, bold=False, align=PP_ALIGN.LEFT,
                 line_spacing=1.4):
        x = x if x is not None else self.M
        y = y if y is not None else Inches(2.5)
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(3.0)
        font_size = font_size or self.theme.body_size
        color = color or self.theme.gray

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(text.split("\n")):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = hex_to_rgbcolor(color)
            p.font.name = self.theme.body_font
            p.alignment = align
            p.space_after = Pt(font_size * (line_spacing - 1))
        return txBox

    def add_bullet_list(self, slide, items: List[str], x=None, y=None, w=None, h=None,
                        font_size=None, color=None, bullet_color=None):
        x = x if x is not None else self.M
        y = y if y is not None else Inches(2.5)
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(3.5)
        font_size = font_size or self.theme.body_size
        color = color or self.theme.gray
        bullet_color = bullet_color or self.theme.day_blue

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run_bullet = p.add_run()
            run_bullet.text = "●  "
            run_bullet.font.size = Pt(font_size - 2)
            run_bullet.font.color.rgb = hex_to_rgbcolor(bullet_color)
            run_bullet.font.name = self.theme.body_font

            run_text = p.add_run()
            run_text.text = item
            run_text.font.size = Pt(font_size)
            run_text.font.color.rgb = hex_to_rgbcolor(color)
            run_text.font.name = self.theme.body_font

            p.space_after = Pt(8)
        return txBox

    # --- Logo helper ---

    def add_logo(self, slide, variant: str = "colored", position: str = "upper-left",
                 max_width_inches: float = 2.0, max_height_inches: float = 0.6):
        allowed = self.theme.brand.logo_rules.get("allowed_placements",
                    ["upper-left", "lower-left", "upper-center", "lower-center"])
        if position not in allowed:
            position = "upper-left"

        logo_map = {
            "colored": self.theme.logo_colored_horizontal,
            "white": self.theme.logo_white_horizontal,
            "black": self.theme.logo_black_horizontal,
            "favicon": self.theme.favicon_colored,
        }
        logo_path = logo_map.get(variant, "")
        if not logo_path or not os.path.exists(logo_path):
            return None

        if Image:
            with Image.open(logo_path) as im:
                nat_w, nat_h = im.size
        else:
            nat_w, nat_h = 1841, 483

        aspect = nat_w / nat_h
        w = Inches(max_width_inches)
        h = Inches(max_width_inches / aspect)
        if h > Inches(max_height_inches):
            h = Inches(max_height_inches)
            w = Inches(max_height_inches * aspect)

        min_w = Inches(1.04)
        if w < min_w:
            w = min_w
            h = Inches(1.04 / aspect)

        M = self.M
        if position == "upper-left":
            x, y = M, M * 0.6
        elif position == "lower-left":
            x, y = M, self.H - h - M * 0.6
        elif position == "upper-center":
            x, y = (self.W - w) / 2, M * 0.6
        elif position == "lower-center":
            x, y = (self.W - w) / 2, self.H - h - M * 0.6
        else:
            x, y = M, M * 0.6

        pic = slide.shapes.add_picture(logo_path, int(x), int(y), int(w), int(h))
        return pic

    # --- Decorative elements ---

    def add_accent_bar(self, slide, x, y, w, h, color=None):
        color = color or self.theme.day_blue
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        set_shape_fill(shape, color)
        set_no_border(shape)
        return shape

    def add_card(self, slide, x, y, w, h, fill_color="#FFFFFF",
                 border_color=None, shadow=True):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        set_shape_fill(shape, fill_color)
        radius_emu = Pt(self.theme.card_radius_pt)
        set_shape_rounded_rect_radius(shape, radius_emu)

        if border_color:
            shape.line.color.rgb = hex_to_rgbcolor(border_color)
            shape.line.width = Pt(1)
        else:
            set_no_border(shape)

        if shadow:
            spPr = get_spPr(shape)
            if spPr is not None:
                nsuri = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                effectLst = etree.SubElement(spPr, f'{{{nsuri}}}effectLst')
                outerShdw = etree.SubElement(effectLst, f'{{{nsuri}}}outerShdw')
                outerShdw.set('blurRad', '152400')
                outerShdw.set('dist', '38100')
                outerShdw.set('dir', '5400000')
                outerShdw.set('rotWithShape', '0')
                srgbClr = etree.SubElement(outerShdw, f'{{{nsuri}}}srgbClr')
                srgbClr.set('val', '000000')
                alpha = etree.SubElement(srgbClr, f'{{{nsuri}}}alpha')
                alpha.set('val', f'{self.theme.card_shadow_alpha * 1000}')
        return shape

    def add_button(self, slide, text: str, x, y, w=None, h=None,
                   fill_color=None, text_color=None):
        fill_color = fill_color or self.theme.button_fill_color
        text_color = text_color or self.theme.button_text_color
        w = w or Inches(2.2)
        h = h or Inches(0.55)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        set_shape_fill(shape, fill_color)
        set_no_border(shape)
        set_shape_rounded_rect_radius(shape, int(h / 2))

        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgbcolor(text_color)
        p.font.name = self.theme.body_font
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        txBody = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')
        return shape

    def add_placeholder_image(self, slide, x, y, w, h, label="Image",
                              fill_color=None):
        fill_color = fill_color or "#E8EDFB"
        shape = self.add_card(slide, x, y, w, h, fill_color=fill_color, shadow=False)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[ {label} ]"
        p.font.size = Pt(12)
        p.font.color.rgb = hex_to_rgbcolor(self.theme.day_blue)
        p.font.name = self.theme.utility_font
        p.alignment = PP_ALIGN.CENTER
        txBody = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')
        return shape

    def add_footer(self, slide, text: str = f"© {_CURRENT_YEAR} OpenTeams  |  openteams.com",
                   show_logo: bool = True, bg_color: str = None):
        footer_h = Inches(0.5)
        y = self.H - footer_h

        if bg_color:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, self.W, footer_h)
            set_shape_fill(bar, bg_color)
            set_no_border(bar)

        text_color = auto_text_color(bg_color or "#FFFFFF")

        txBox = slide.shapes.add_textbox(self.M, y + Inches(0.08), Inches(8), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(self.theme.caption_size)
        p.font.color.rgb = hex_to_rgbcolor(text_color)
        p.font.name = self.theme.utility_font
        p.alignment = PP_ALIGN.LEFT

        if show_logo and self.theme.favicon_colored and os.path.exists(self.theme.favicon_colored):
            icon_size = Inches(0.3)
            slide.shapes.add_picture(
                self.theme.favicon_colored,
                int(self.W - self.M - icon_size),
                int(y + (footer_h - icon_size) / 2),
                int(icon_size), int(icon_size)
            )

    def add_section_header(self, slide, title: str, subtitle: str = "",
                           bg_color: str = None):
        bg_color = bg_color or self.theme.night_navy
        add_slide_bg_color(slide, bg_color)
        text_color = auto_text_color(bg_color)

        self.add_accent_bar(slide, self.M, Inches(2.8), Inches(0.8), Inches(0.06),
                           color=self.theme.day_blue if bg_color != self.theme.day_blue else self.theme.yellow)

        self.add_title(slide, title, y=Inches(3.0), font_size=self.theme.h1_size,
                      color=text_color)

        if subtitle:
            sub_color = self.theme.day_blue if luminance(bg_color) < 0.3 else self.theme.night_navy
            if contrast_ratio(bg_color, sub_color) < 3:
                sub_color = text_color
            self.add_body(slide, subtitle, y=Inches(4.2), font_size=self.theme.h3_size,
                         color=sub_color)

    def add_metric_card(self, slide, x, y, w, h, value: str, label: str,
                        accent_color: str = None):
        accent_color = accent_color or self.theme.day_blue
        self.add_card(slide, x, y, w, h)

        bar_h = Inches(0.06)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, bar_h)
        set_shape_fill(bar, accent_color)
        set_no_border(bar)

        val_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3),
                                            w - Inches(0.4), Inches(0.8))
        tf = val_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgbcolor(self.theme.night_navy)
        p.font.name = self.theme.headline_font
        p.alignment = PP_ALIGN.LEFT

        lbl_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.0),
                                            w - Inches(0.4), Inches(0.5))
        tf2 = lbl_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = hex_to_rgbcolor(self.theme.gray)
        p2.font.name = self.theme.body_font
        p2.alignment = PP_ALIGN.LEFT

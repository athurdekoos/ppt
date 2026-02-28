"""
Slide renderers — one function per slide type.
Each renderer signature: render_*(sb: SlideBuilder, spec: dict) -> None

The SlideBuilder class provides high-level branded helpers.
The spec dict comes from the user's slide JSON.
"""
from __future__ import annotations

import os
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

from pptx_helpers import (
    hex_to_rgbcolor, luminance, contrast_ratio, auto_text_color,
    set_shape_fill, set_shape_rounded_rect_radius, set_no_border,
    add_slide_bg_color, make_gradient_rect, set_shape_alpha, _get_spPr,
)
from brand_engine import ThemeConfig

try:
    from PIL import Image
except ImportError:
    Image = None


# ===================================================================
# SlideBuilder — high-level helpers for building branded slides
# ===================================================================

class SlideBuilder:
    """High-level helpers for building branded slides."""

    def __init__(self, prs: Presentation, theme: ThemeConfig):
        self.prs = prs
        self.theme = theme
        self.c = theme  # color shortcuts via theme
        self.t = theme  # typography shortcuts via theme
        self.W = Inches(theme.slide_width_inches)
        self.H = Inches(theme.slide_height_inches)
        self.M = Inches(theme.margin_inches)
        self.G = Inches(theme.gutter_inches)

    def new_slide(self):
        """Add a blank slide."""
        layout = self.prs.slide_layouts[6]  # blank layout
        return self.prs.slides.add_slide(layout)

    # --- Text helpers ---

    def add_title(self, slide, text: str, x=None, y=None, w=None, h=None,
                  font_size=None, color=None, bold=True, align=PP_ALIGN.LEFT):
        x = x if x is not None else self.M
        y = y if y is not None else self.M
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(1.0)
        font_size = font_size or self.theme.h1_size
        color = color or self.theme.night_navy

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = hex_to_rgbcolor(color)
        p.font.name = self.theme.headline_font
        p.alignment = align
        return txBox

    def add_subtitle(self, slide, text: str, x=None, y=None, w=None, h=None,
                     font_size=None, color=None):
        x = x if x is not None else self.M
        y = y if y is not None else Inches(1.6)
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(0.8)
        font_size = font_size or self.theme.h3_size
        color = color or self.theme.day_blue

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = False
        p.font.color.rgb = hex_to_rgbcolor(color)
        p.font.name = self.theme.body_font
        p.alignment = PP_ALIGN.LEFT
        return txBox

    def add_body(self, slide, text: str, x=None, y=None, w=None, h=None,
                 font_size=None, color=None, bold=False, align=PP_ALIGN.LEFT,
                 line_spacing=1.4):
        x = x if x is not None else self.M
        y = y if y is not None else Inches(2.5)
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(3.0)
        font_size = font_size or self.theme.body_size
        color = color or self.theme.gray

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(text.split("\n")):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = hex_to_rgbcolor(color)
            p.font.name = self.theme.body_font
            p.alignment = align
            p.space_after = Pt(font_size * (line_spacing - 1))
        return txBox

    def add_bullet_list(self, slide, items: List[str], x=None, y=None, w=None, h=None,
                        font_size=None, color=None, bullet_color=None):
        x = x if x is not None else self.M
        y = y if y is not None else Inches(2.5)
        w = w if w is not None else (self.W - 2 * self.M)
        h = h if h is not None else Inches(3.5)
        font_size = font_size or self.theme.body_size
        color = color or self.theme.gray
        bullet_color = bullet_color or self.theme.day_blue

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run_bullet = p.add_run()
            run_bullet.text = "●  "
            run_bullet.font.size = Pt(font_size - 2)
            run_bullet.font.color.rgb = hex_to_rgbcolor(bullet_color)
            run_bullet.font.name = self.theme.body_font

            run_text = p.add_run()
            run_text.text = item
            run_text.font.size = Pt(font_size)
            run_text.font.color.rgb = hex_to_rgbcolor(color)
            run_text.font.name = self.theme.body_font

            p.space_after = Pt(8)
        return txBox

    # --- Logo helper ---

    def add_logo(self, slide, variant: str = "colored", position: str = "upper-left",
                 max_width_inches: float = 2.0, max_height_inches: float = 0.6):
        allowed = self.theme.brand.logo_rules.get("allowed_placements",
                    ["upper-left", "lower-left", "upper-center", "lower-center"])
        if position not in allowed:
            position = "upper-left"

        logo_map = {
            "colored": self.theme.logo_colored_horizontal,
            "white": self.theme.logo_white_horizontal,
            "black": self.theme.logo_black_horizontal,
            "favicon": self.theme.favicon_colored,
        }
        logo_path = logo_map.get(variant, "")
        if not logo_path or not os.path.exists(logo_path):
            return None

        if Image:
            with Image.open(logo_path) as im:
                nat_w, nat_h = im.size
        else:
            nat_w, nat_h = 1841, 483

        aspect = nat_w / nat_h
        w = Inches(max_width_inches)
        h = Inches(max_width_inches / aspect)
        if h > Inches(max_height_inches):
            h = Inches(max_height_inches)
            w = Inches(max_height_inches * aspect)

        min_w = Inches(1.04)
        if w < min_w:
            w = min_w
            h = Inches(1.04 / aspect)

        M = self.M
        if position == "upper-left":
            x, y = M, M * 0.6
        elif position == "lower-left":
            x, y = M, self.H - h - M * 0.6
        elif position == "upper-center":
            x, y = (self.W - w) / 2, M * 0.6
        elif position == "lower-center":
            x, y = (self.W - w) / 2, self.H - h - M * 0.6
        else:
            x, y = M, M * 0.6

        pic = slide.shapes.add_picture(logo_path, int(x), int(y), int(w), int(h))
        return pic

    # --- Decorative elements ---

    def add_accent_bar(self, slide, x, y, w, h, color=None):
        color = color or self.theme.day_blue
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        set_shape_fill(shape, color)
        set_no_border(shape)
        return shape

    def add_card(self, slide, x, y, w, h, fill_color="#FFFFFF",
                 border_color=None, shadow=True):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        set_shape_fill(shape, fill_color)
        radius_emu = Pt(self.theme.card_radius_pt)
        set_shape_rounded_rect_radius(shape, radius_emu)

        if border_color:
            shape.line.color.rgb = hex_to_rgbcolor(border_color)
            shape.line.width = Pt(1)
        else:
            set_no_border(shape)

        if shadow:
            spPr = _get_spPr(shape)
            if spPr is not None:
                nsuri = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                effectLst = etree.SubElement(spPr, f'{{{nsuri}}}effectLst')
                outerShdw = etree.SubElement(effectLst, f'{{{nsuri}}}outerShdw')
                outerShdw.set('blurRad', '152400')
                outerShdw.set('dist', '38100')
                outerShdw.set('dir', '5400000')
                outerShdw.set('rotWithShape', '0')
                srgbClr = etree.SubElement(outerShdw, f'{{{nsuri}}}srgbClr')
                srgbClr.set('val', '000000')
                alpha = etree.SubElement(srgbClr, f'{{{nsuri}}}alpha')
                alpha.set('val', f'{self.theme.card_shadow_alpha * 1000}')
        return shape

    def add_button(self, slide, text: str, x, y, w=None, h=None,
                   fill_color=None, text_color=None):
        fill_color = fill_color or self.theme.button_fill_color
        text_color = text_color or self.theme.button_text_color
        w = w or Inches(2.2)
        h = h or Inches(0.55)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        set_shape_fill(shape, fill_color)
        set_no_border(shape)
        set_shape_rounded_rect_radius(shape, int(h / 2))

        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgbcolor(text_color)
        p.font.name = self.theme.body_font
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        txBody = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')
        return shape

    def add_placeholder_image(self, slide, x, y, w, h, label="Image",
                              fill_color=None):
        fill_color = fill_color or "#E8EDFB"
        shape = self.add_card(slide, x, y, w, h, fill_color=fill_color, shadow=False)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[ {label} ]"
        p.font.size = Pt(12)
        p.font.color.rgb = hex_to_rgbcolor(self.theme.day_blue)
        p.font.name = self.theme.utility_font
        p.alignment = PP_ALIGN.CENTER
        txBody = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')
        return shape

    def add_footer(self, slide, text: str = "© 2025 OpenTeams  |  openteams.com",
                   show_logo: bool = True, bg_color: str = None):
        footer_h = Inches(0.5)
        y = self.H - footer_h

        if bg_color:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, self.W, footer_h)
            set_shape_fill(bar, bg_color)
            set_no_border(bar)

        text_color = auto_text_color(bg_color or "#FFFFFF")

        txBox = slide.shapes.add_textbox(self.M, y + Inches(0.08), Inches(8), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(self.theme.caption_size)
        p.font.color.rgb = hex_to_rgbcolor(text_color)
        p.font.name = self.theme.utility_font
        p.alignment = PP_ALIGN.LEFT

        if show_logo and self.theme.favicon_colored and os.path.exists(self.theme.favicon_colored):
            icon_size = Inches(0.3)
            slide.shapes.add_picture(
                self.theme.favicon_colored,
                int(self.W - self.M - icon_size),
                int(y + (footer_h - icon_size) / 2),
                int(icon_size), int(icon_size)
            )

    def add_section_header(self, slide, title: str, subtitle: str = "",
                           bg_color: str = None):
        bg_color = bg_color or self.theme.night_navy
        add_slide_bg_color(slide, bg_color)
        text_color = auto_text_color(bg_color)

        self.add_accent_bar(slide, self.M, Inches(2.8), Inches(0.8), Inches(0.06),
                           color=self.theme.day_blue if bg_color != self.theme.day_blue else self.theme.yellow)

        self.add_title(slide, title, y=Inches(3.0), font_size=self.theme.h1_size,
                      color=text_color)

        if subtitle:
            sub_color = self.theme.day_blue if luminance(bg_color) < 0.3 else self.theme.night_navy
            if contrast_ratio(bg_color, sub_color) < 3:
                sub_color = text_color
            self.add_body(slide, subtitle, y=Inches(4.2), font_size=self.theme.h3_size,
                         color=sub_color)

    def add_metric_card(self, slide, x, y, w, h, value: str, label: str,
                        accent_color: str = None):
        accent_color = accent_color or self.theme.day_blue
        self.add_card(slide, x, y, w, h)

        bar_h = Inches(0.06)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, bar_h)
        set_shape_fill(bar, accent_color)
        set_no_border(bar)

        val_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3),
                                            w - Inches(0.4), Inches(0.8))
        tf = val_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgbcolor(self.theme.night_navy)
        p.font.name = self.theme.headline_font
        p.alignment = PP_ALIGN.LEFT

        lbl_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.0),
                                            w - Inches(0.4), Inches(0.5))
        tf2 = lbl_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = hex_to_rgbcolor(self.theme.gray)
        p2.font.name = self.theme.body_font
        p2.alignment = PP_ALIGN.LEFT


# ===================================================================
# Slide Renderers — one per slide type
# ===================================================================

def render_cover(sb: SlideBuilder, spec: dict) -> None:
    """Cover / Title slide (hero layout)."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")

    # Gradient accent block on right
    make_gradient_rect(slide, Inches(7.5), 0, Inches(5.833), sb.H,
                       sb.theme.night_navy, sb.theme.day_blue, angle=135)

    # Decorative dots
    for dx, dy, col in [
        (8.2, 1.5, sb.theme.yellow), (9.8, 2.0, sb.theme.salmon),
        (8.8, 5.0, sb.theme.day_blue), (10.5, 4.2, sb.theme.yellow),
    ]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx), Inches(dy),
                                      Inches(0.35), Inches(0.35))
        set_shape_fill(dot, col)
        set_no_border(dot)
        set_shape_alpha(dot, '30000')

    # Decorative favicon on gradient panel
    if sb.theme.favicon_colored and os.path.exists(sb.theme.favicon_colored):
        fav_size = Inches(3.5)
        slide.shapes.add_picture(
            sb.theme.favicon_colored,
            int(Inches(9.0) - fav_size / 2), int(Inches(3.75) - fav_size / 2),
            int(fav_size), int(fav_size)
        )

    # Logo
    sb.add_logo(slide, "colored", "upper-left", max_width_inches=2.4, max_height_inches=0.65)

    # Title
    title = spec.get("title", "Presentation Title")
    sb.add_title(slide, title,
                x=sb.M, y=Inches(2.4), w=Inches(6.5), h=Inches(1.8),
                font_size=48, color=sb.theme.night_navy)

    # Subtitle + date
    subtitle = spec.get("subtitle", "")
    date = spec.get("date", "")
    sub_text = subtitle
    if date:
        sub_text = f"{subtitle}\n{date}" if subtitle else date
    if sub_text:
        sb.add_body(slide, sub_text,
                   x=sb.M, y=Inches(4.3), w=Inches(6.0), h=Inches(1.2),
                   font_size=20, color=sb.theme.day_blue)

    # CTA button
    sb.add_button(slide, "Get Started", sb.M, Inches(5.8))

    # Footer
    sb.add_footer(slide, "Confidential  |  © 2025 OpenTeams", show_logo=False)


def render_section_divider(sb: SlideBuilder, spec: dict) -> None:
    """Section Divider slide."""
    slide = sb.new_slide()
    title = spec.get("title", "Section Title")
    subtitle = spec.get("subtitle", "")
    bg_color = spec.get("bg_color", None)
    sb.add_section_header(slide, title, subtitle, bg_color=bg_color)
    sb.add_logo(slide, "white", "lower-left", max_width_inches=1.8, max_height_inches=0.45)


def render_agenda(sb: SlideBuilder, spec: dict) -> None:
    """Agenda slide."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")

    sb.add_accent_bar(slide, sb.M, Inches(0.9), Inches(0.8), Inches(0.06))
    sb.add_title(slide, "Agenda", y=Inches(1.1), font_size=sb.theme.h2_size)

    items = spec.get("items", ["Topic 1", "Topic 2", "Topic 3"])
    accent_colors = [sb.theme.day_blue, sb.theme.night_navy, sb.theme.yellow,
                     sb.theme.salmon, sb.theme.day_blue]

    y_start = Inches(2.4)
    for i, item in enumerate(items):
        y = y_start + Inches(i * 0.85)
        accent = accent_colors[i % len(accent_colors)]

        # Number circle
        num_size = Inches(0.5)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, sb.M, y, num_size, num_size)
        set_shape_fill(circle, accent)
        set_no_border(circle)
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgbcolor(auto_text_color(accent))
        p.font.name = sb.theme.headline_font
        p.alignment = PP_ALIGN.CENTER
        txBody = circle._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')

        sb.add_body(slide, item, x=sb.M + Inches(0.75), y=y + Inches(0.05),
                   w=Inches(10), h=Inches(0.5), font_size=18, color=sb.theme.gray, bold=False)

    sb.add_footer(slide, show_logo=True)
    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)


def render_content(sb: SlideBuilder, spec: dict) -> None:
    """Content slide (title + body + visual placeholder)."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")

    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)
    sb.add_accent_bar(slide, sb.M, Inches(1.2), Inches(0.8), Inches(0.06))

    title = spec.get("title", "Content Slide Title")
    sb.add_title(slide, title, y=Inches(1.4), font_size=sb.theme.h2_size)

    # Body text or bullet list on left
    body = spec.get("body", "")
    bullet_items = spec.get("bullet_items", [])
    if bullet_items:
        sb.add_bullet_list(slide, bullet_items,
                          x=sb.M, y=Inches(2.5), w=Inches(5.5), h=Inches(3.5),
                          font_size=15)
    elif body:
        sb.add_body(slide, body,
                   x=sb.M, y=Inches(2.5), w=Inches(5.5), h=Inches(3.5),
                   font_size=15, color=sb.theme.gray)
    else:
        sb.add_body(slide, "Add your key points here.",
                   x=sb.M, y=Inches(2.5), w=Inches(5.5), h=Inches(3.5),
                   font_size=15, color=sb.theme.gray)

    # Image placeholder on right
    img_label = spec.get("image_placeholder", "Visual / Image")
    sb.add_placeholder_image(slide, Inches(7.0), Inches(1.4),
                             Inches(5.7), Inches(4.8), img_label)

    sb.add_footer(slide, show_logo=True)


def render_two_column(sb: SlideBuilder, spec: dict) -> None:
    """Two-column content slide."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")

    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)
    sb.add_accent_bar(slide, sb.M, Inches(1.2), Inches(0.8), Inches(0.06))

    title = spec.get("title", "Two-Column Layout")
    sb.add_title(slide, title, y=Inches(1.4), font_size=sb.theme.h2_size)

    col_w = Inches(5.8)
    col1_x = sb.M
    col2_x = sb.M + col_w + sb.G

    # Column 1
    sb.add_card(slide, col1_x, Inches(2.5), col_w, Inches(4.0))
    left_title = spec.get("left_title", "Left Column")
    left_body = spec.get("left_body", "")
    sb.add_body(slide, left_title, x=col1_x + Inches(0.3), y=Inches(2.7),
               w=col_w - Inches(0.6), h=Inches(0.5),
               font_size=20, color=sb.theme.night_navy, bold=True)
    if left_body:
        sb.add_body(slide, left_body, x=col1_x + Inches(0.3), y=Inches(3.3),
                   w=col_w - Inches(0.6), h=Inches(2.5),
                   font_size=14, color=sb.theme.gray)

    # Column 2
    sb.add_card(slide, col2_x, Inches(2.5), col_w, Inches(4.0))
    right_title = spec.get("right_title", "Right Column")
    right_body = spec.get("right_body", "")
    sb.add_body(slide, right_title, x=col2_x + Inches(0.3), y=Inches(2.7),
               w=col_w - Inches(0.6), h=Inches(0.5),
               font_size=20, color=sb.theme.night_navy, bold=True)
    if right_body:
        sb.add_body(slide, right_body, x=col2_x + Inches(0.3), y=Inches(3.3),
                   w=col_w - Inches(0.6), h=Inches(2.5),
                   font_size=14, color=sb.theme.gray)

    sb.add_footer(slide, show_logo=True)


def render_quote(sb: SlideBuilder, spec: dict) -> None:
    """Big statement / quote slide."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, sb.theme.night_navy)

    # Large decorative quote mark
    quote_mark = slide.shapes.add_textbox(sb.M, Inches(1.0), Inches(2), Inches(2))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    p.text = "\u201C"
    p.font.size = Pt(160)
    p.font.color.rgb = hex_to_rgbcolor(sb.theme.day_blue)
    p.font.name = sb.theme.headline_font
    p.font.bold = True

    # Quote text
    text = spec.get("text", "A bold statement that captures\nyour key message in one line.")
    sb.add_title(slide, text,
                x=Inches(1.2), y=Inches(2.8), w=Inches(10.5), h=Inches(2.0),
                font_size=36, color="#FFFFFF")

    # Attribution
    attribution = spec.get("attribution", "")
    if attribution:
        sb.add_body(slide, f"— {attribution}",
                   x=Inches(1.2), y=Inches(5.0), w=Inches(8), h=Inches(0.6),
                   font_size=16, color=sb.theme.day_blue)

    # Accent dots
    for dx, dy, col in [(11.5, 5.5, sb.theme.yellow), (12.0, 5.0, sb.theme.salmon)]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx), Inches(dy),
                                      Inches(0.25), Inches(0.25))
        set_shape_fill(dot, col)
        set_no_border(dot)

    sb.add_logo(slide, "white", "lower-left", max_width_inches=1.8, max_height_inches=0.45)


def render_metrics(sb: SlideBuilder, spec: dict) -> None:
    """Data/metrics slide with metric cards."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#F7F8FC")

    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)
    sb.add_accent_bar(slide, sb.M, Inches(1.1), Inches(0.8), Inches(0.06))

    title = spec.get("title", "Key Metrics")
    sb.add_title(slide, title, y=Inches(1.3), font_size=sb.theme.h2_size)

    metrics = spec.get("metrics", [
        {"value": "98%", "label": "Metric 1"},
        {"value": "3.5x", "label": "Metric 2"},
        {"value": "500+", "label": "Metric 3"},
        {"value": "24/7", "label": "Metric 4"},
    ])

    accent_colors = [sb.theme.day_blue, sb.theme.night_navy, sb.theme.yellow, sb.theme.salmon]

    num_metrics = len(metrics)
    if num_metrics == 0:
        return

    # Calculate card layout
    total_gap = Inches(0.35) * (num_metrics - 1)
    usable_w = sb.W - 2 * sb.M - total_gap
    card_w_val = min(usable_w / num_metrics, Inches(2.7))
    card_h = Inches(1.6)
    gap = Inches(0.35)
    start_x = sb.M
    y = Inches(2.3)

    for i, metric in enumerate(metrics):
        x = start_x + i * (card_w_val + gap)
        accent = accent_colors[i % len(accent_colors)]
        sb.add_metric_card(slide, x, y, card_w_val, card_h,
                          metric.get("value", "—"), metric.get("label", ""),
                          accent)

    # Chart placeholder below
    sb.add_placeholder_image(slide, sb.M, Inches(4.3), Inches(11.8), Inches(2.5),
                             "Chart / Data Visualization")

    sb.add_footer(slide, show_logo=True)


def render_team(sb: SlideBuilder, spec: dict) -> None:
    """Team / Profile slide."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")

    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)
    sb.add_accent_bar(slide, sb.M, Inches(1.1), Inches(0.8), Inches(0.06))

    title = spec.get("title", "Our Team")
    sb.add_title(slide, title, y=Inches(1.3), font_size=sb.theme.h2_size)

    members = spec.get("members", [
        {"name": f"Team Member {i+1}", "role": "Role / Title", "bio": "Brief bio."}
        for i in range(4)
    ])

    accent_colors = [sb.theme.day_blue, sb.theme.night_navy, sb.theme.yellow, sb.theme.salmon]

    num_members = min(len(members), 6)  # cap at 6
    card_w = Inches(2.7)
    card_h = Inches(4.2)
    gap = Inches(0.35)
    start_x = sb.M + Inches(0.3)
    y = Inches(2.3)

    for i in range(num_members):
        member = members[i]
        x = start_x + i * (card_w + gap)
        accent = accent_colors[i % len(accent_colors)]

        sb.add_card(slide, x, y, card_w, card_h)

        # Avatar circle
        avatar_size = Inches(1.4)
        avatar_x = x + (card_w - avatar_size) / 2
        avatar_y = y + Inches(0.4)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(avatar_x), int(avatar_y),
                                         int(avatar_size), int(avatar_size))
        set_shape_fill(circle, accent)
        set_no_border(circle)
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = "👤"
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER
        txBody = circle._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')

        # Name
        name = member.get("name", f"Team Member {i+1}")
        sb.add_body(slide, name,
                   x=x + Inches(0.2), y=y + Inches(2.1),
                   w=card_w - Inches(0.4), h=Inches(0.4),
                   font_size=16, color=sb.theme.night_navy, bold=True, align=PP_ALIGN.CENTER)

        # Role
        role = member.get("role", "Role / Title")
        sb.add_body(slide, role,
                   x=x + Inches(0.2), y=y + Inches(2.6),
                   w=card_w - Inches(0.4), h=Inches(0.3),
                   font_size=12, color=sb.theme.day_blue, align=PP_ALIGN.CENTER)

        # Bio
        bio = member.get("bio", "")
        if bio:
            sb.add_body(slide, bio,
                       x=x + Inches(0.2), y=y + Inches(3.1),
                       w=card_w - Inches(0.4), h=Inches(0.8),
                       font_size=11, color=sb.theme.gray, align=PP_ALIGN.CENTER)

    sb.add_footer(slide, show_logo=True)


def render_case_study(sb: SlideBuilder, spec: dict) -> None:
    """Case Study (Challenge → Solution → Results)."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")

    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)
    sb.add_accent_bar(slide, sb.M, Inches(1.1), Inches(0.8), Inches(0.06))

    title = spec.get("title", "Case Study: Client Name")
    sb.add_title(slide, title, y=Inches(1.3), font_size=sb.theme.h2_size)

    col_w = Inches(3.7)
    gap = Inches(0.4)
    start_x = sb.M + Inches(0.15)
    y = Inches(2.5)
    col_h = Inches(4.0)

    labels = ["Challenge", "Solution", "Results"]
    colors = [sb.theme.salmon, sb.theme.day_blue, sb.theme.yellow]
    icons = ["⚡", "🔧", "📈"]
    bodies = [
        spec.get("challenge", "Describe the challenge."),
        spec.get("solution", "Describe the solution."),
        spec.get("results", "Describe the results."),
    ]

    for i, (label, accent, icon, body) in enumerate(zip(labels, colors, icons, bodies)):
        x = start_x + i * (col_w + gap)

        sb.add_card(slide, x, y, col_w, col_h)

        # Accent top bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_w, Inches(0.07))
        set_shape_fill(bar, accent)
        set_no_border(bar)

        # Icon circle
        icon_size = Inches(0.6)
        icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       int(x + Inches(0.3)), int(y + Inches(0.4)),
                                       int(icon_size), int(icon_size))
        set_shape_fill(icon_shape, accent)
        set_no_border(icon_shape)
        tf = icon_shape.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
        txBody = icon_shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if txBody is not None:
            txBody.set('anchor', 'ctr')

        # Label
        sb.add_body(slide, label,
                   x=x + Inches(0.3), y=y + Inches(1.2),
                   w=col_w - Inches(0.6), h=Inches(0.4),
                   font_size=20, color=sb.theme.night_navy, bold=True)

        # Body
        sb.add_body(slide, body,
                   x=x + Inches(0.3), y=y + Inches(1.8),
                   w=col_w - Inches(0.6), h=Inches(1.8),
                   font_size=13, color=sb.theme.gray)

    sb.add_footer(slide, show_logo=True)


def render_closing(sb: SlideBuilder, spec: dict) -> None:
    """Closing / CTA slide."""
    slide = sb.new_slide()

    # Gradient background
    make_gradient_rect(slide, 0, 0, sb.W, sb.H, sb.theme.night_navy, sb.theme.day_blue, angle=135)

    # Decorative circles
    for dx, dy, sz, col, alpha_val in [
        (1.5, 1.0, 1.5, sb.theme.yellow, "15000"),
        (10.5, 5.5, 2.0, sb.theme.salmon, "12000"),
        (11.0, 1.5, 0.8, sb.theme.day_blue, "20000"),
    ]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx), Inches(dy),
                                      Inches(sz), Inches(sz))
        set_shape_fill(dot, col)
        set_no_border(dot)
        set_shape_alpha(dot, alpha_val)

    # Title
    title = spec.get("title", "Thank You")
    sb.add_title(slide, title,
                x=sb.M, y=Inches(2.0), w=Inches(12), h=Inches(1.5),
                font_size=56, color="#FFFFFF", align=PP_ALIGN.CENTER)

    # Subtitle
    subtitle = spec.get("subtitle", "Questions? Let's discuss.")
    if subtitle:
        sb.add_body(slide, subtitle,
                   x=sb.M, y=Inches(3.6), w=Inches(12), h=Inches(0.8),
                   font_size=22, color=sb.theme.day_blue, align=PP_ALIGN.CENTER)

    # CTA button
    cta_text = spec.get("cta_text", "Contact Us")
    btn_w = Inches(2.8)
    sb.add_button(slide, cta_text,
                 (sb.W - btn_w) / 2, Inches(4.8), w=btn_w,
                 fill_color="#FFFFFF", text_color=sb.theme.night_navy)

    # Contact info
    contact = spec.get("contact", "hello@openteams.com  |  openteams.com")
    if contact:
        sb.add_body(slide, contact,
                   x=sb.M, y=Inches(5.8), w=Inches(12), h=Inches(0.5),
                   font_size=14, color="#FFFFFF", align=PP_ALIGN.CENTER)

    # Logo
    sb.add_logo(slide, "white", "lower-center", max_width_inches=2.2, max_height_inches=0.55)


def render_blank(sb: SlideBuilder, spec: dict) -> None:
    """Blank slide with logo only."""
    slide = sb.new_slide()
    add_slide_bg_color(slide, "#FFFFFF")
    sb.add_logo(slide, "colored", "upper-left", max_width_inches=1.8, max_height_inches=0.45)


# ===================================================================
# Registry
# ===================================================================

RENDERERS = {
    "cover":            render_cover,
    "section_divider":  render_section_divider,
    "agenda":           render_agenda,
    "content":          render_content,
    "two_column":       render_two_column,
    "quote":            render_quote,
    "metrics":          render_metrics,
    "team":             render_team,
    "case_study":       render_case_study,
    "closing":          render_closing,
    "blank":            render_blank,
}

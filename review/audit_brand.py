#!/usr/bin/env python3
"""Audit a .pptx file against OpenTeams 2025 brand rules."""
import sys
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.dml import MSO_THEME_COLOR

ALLOWED_COLORS_HEX = {
    "022791", "4D75FE", "FF8A69", "FAA944",
    "0C0C0C", "262626", "3F3F3F", "FFFFFF",
    "F7F8FC", "E8EDFB", "3AD58E",
    "000000",  # default black
}

ACCENT_ONLY = {"FF8A69", "FAA944"}
ALLOWED_FONTS = {"Inter Tight", "Roboto", "Arial", None}

def hex_of(rgb):
    if rgb is None:
        return None
    return str(rgb).upper()

def audit(path):
    prs = Presentation(path)
    results = []

    # C18: Slide dimensions
    w = prs.slide_width / 914400
    h = prs.slide_height / 914400
    if abs(w - 13.333) > 0.05 or abs(h - 7.5) > 0.05:
        results.append(f"C18 FAIL: Slide size {w:.3f}x{h:.3f} (expected 13.333x7.5)")
    else:
        results.append(f"C18 PASS: Slide size {w:.3f}x{h:.3f}")

    total_color_checks = 0
    total_font_checks = 0
    color_issues = []
    font_issues = []
    logo_info = []

    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            # Check fill colors
            try:
                fill = shape.fill
                if fill.type is not None:
                    try:
                        c = hex_of(fill.fore_color.rgb)
                        total_color_checks += 1
                        if c and c not in ALLOWED_COLORS_HEX:
                            color_issues.append(f"  Slide {i}, shape '{shape.name}': fill #{c} NOT in brand palette")
                    except (AttributeError, TypeError):
                        pass
            except (AttributeError, TypeError):
                pass

            # Check text colors and fonts
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Font color
                        try:
                            if run.font.color and run.font.color.rgb:
                                c = hex_of(run.font.color.rgb)
                                total_color_checks += 1
                                if c and c not in ALLOWED_COLORS_HEX:
                                    color_issues.append(
                                        f"  Slide {i}, text '{run.text[:40]}': color #{c} NOT in brand palette")
                        except (AttributeError, TypeError):
                            pass

                        # Font name (C5-C8)
                        fname = run.font.name
                        if fname is not None:
                            total_font_checks += 1
                            if fname not in ALLOWED_FONTS:
                                font_issues.append(
                                    f"  C8 FAIL Slide {i}: font '{fname}' on text '{run.text[:40]}'")

            # Logo tracking
            if hasattr(shape, 'image'):
                logo_info.append(
                    f"  Slide {i}: Image '{shape.name}' at left={shape.left/914400:.2f}in, "
                    f"top={shape.top/914400:.2f}in, w={shape.width/914400:.2f}in, h={shape.height/914400:.2f}in")

    # Print results
    print("=" * 60)
    print("OpenTeams Brand Compliance Audit")
    print("=" * 60)

    print(f"\n{results[0]}")

    print(f"\n--- COLOR CHECKS ({total_color_checks} inspected) ---")
    if color_issues:
        for ci in color_issues:
            print(ci)
    else:
        print("  ✅ All fill/text colors are within the brand palette")

    print(f"\n--- FONT CHECKS ({total_font_checks} inspected) ---")
    if font_issues:
        for fi in font_issues:
            print(fi)
    else:
        print("  ✅ All fonts are Inter Tight, Roboto, or Arial")

    print(f"\n--- LOGO/IMAGE PLACEMENTS ({len(logo_info)} found) ---")
    for li in logo_info:
        print(li)
    # Check C9: no right-aligned logos
    for li in logo_info:
        # Parse left position - if > 10in on a 13.33in slide, it's right-aligned
        parts = li.split("left=")[1].split("in")[0]
        left_pos = float(parts)
        if left_pos > 10.0:
            print(f"  ⚠️  C9 WARNING: Image at left={left_pos:.2f}in may be right-aligned")

    total_issues = len(color_issues) + len(font_issues)
    print(f"\n{'=' * 60}")
    print(f"Summary: {total_color_checks + total_font_checks + 1} checks | {total_issues} issues")
    if total_issues == 0:
        print("✅ All automated checks PASSED")
    else:
        print(f"❌ {total_issues} issue(s) found — review above")
    print("=" * 60)

    return total_issues

if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else "/home/mia/dev/ppt/review/mock_demo.pptx"
    sys.exit(audit(path))
